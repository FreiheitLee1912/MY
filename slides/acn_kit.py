# -*- coding: utf-8 -*-
"""Shared page system for the slide builders in this folder.

The grid and type scale are lifted from the Accenture 2020 starter pack:
0.42" side margins over a 12.50" text column, an Arial Black headline set flush
top-left with no rule under it, footer furniture on the 7.10" baseline. Only the
palette changes between brands.
"""
import os
import re
import shutil
import tempfile
import zipfile

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def C(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


BLACK, WHITE = C("000000"), C("FFFFFF")
INK = C("333333")

BRANDS = {
    # Accenture 2020 theme: accent1 A100FF, accent2 7500C0, accent3 460073,
    # accent5 BE82FF, dk2 96968C, lt2 E6E6DC
    "acn": dict(
        key=C("A100FF"), mid=C("7500C0"), dark=C("460073"), light=C("BE82FF"),
        tint=C("F4E9FF"), muted=C("70706A"),
        out_fill=C("E6E6DC"), out_fg=C("70706A"),
        pre_fill=C("F4E9FF"), pre_fg=C("460073"), pre_line=C("A100FF"),
        post_fills=(C("A100FF"), C("7500C0")), post_fg=WHITE,
        after_sop=C("460073"),
        open_fill=C("E6E6DC"), open_fg=C("460073"),
        sowhat_fill=C("A100FF"), sowhat_fg=WHITE,
        axis=C("96968C"), track=C("96968C"),
        rule=C("D9D9D2"), panel=C("F7F7F3"),
    ),
    # DAICEL: corporate blue, keeping the source slides' green for "after SOP"
    "daicel": dict(
        key=C("0068B7"), mid=C("1B6EC2"), dark=C("1F4E79"), light=C("9DC3E6"),
        tint=C("E8F1FA"), muted=C("6B7280"),
        out_fill=C("E7E7E7"), out_fg=C("767676"),
        pre_fill=C("E8F1FA"), pre_fg=C("1F4E79"), pre_line=C("0068B7"),
        post_fills=(C("4C9A2A"), C("3D7C21")), post_fg=WHITE,
        after_sop=C("3D7C21"),
        open_fill=C("FBF4E0"), open_fg=C("0068B7"),
        sowhat_fill=C("0068B7"), sowhat_fg=WHITE,
        axis=C("A6A6A6"), track=C("A6A6A6"),
        rule=C("DCDCDC"), panel=C("F5F7FA"),
    ),
}

# semantic colours, shared by both brands
OK, WARN, DENY = C("2E7D32"), C("B7791F"), C("B5292F")
OK_BG, WARN_BG, DENY_BG = C("EAF3EA"), C("FBF4E0"), C("FBEFEE")

JP, LAT, DISPLAY = "Meiryo", "Arial", "Arial Black"

# page grid
L, R = 0.42, 12.92
W = R - L
FOOT_Y = 7.10


# ---------------------------------------------------------------- helpers ----
def style(run, size, bold=False, color=BLACK, face=LAT, ea=JP):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, face
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag, val in (("a:ea", ea), ("a:cs", face)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", val)


def textbox(sl, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tb, tf


def say(sl, x, y, w, h, text, size, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, face=LAT, anchor=MSO_ANCHOR.TOP, space=None):
    _, tf = textbox(sl, x, y, w, h, align, anchor)
    r = tf.paragraphs[0].add_run()
    r.text = text
    style(r, size, bold, color, face)
    if space is not None:
        tf.paragraphs[0].line_spacing = space
    return tf


def rich(sl, x, y, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
    """parts = [(text, size, bold, color) | (text, size, bold, color, face), ...]"""
    _, tf = textbox(sl, x, y, w, h, align, anchor)
    p = tf.paragraphs[0]
    for part in parts:
        text, size, bold, color = part[:4]
        face = part[4] if len(part) > 4 else LAT
        r = p.add_run()
        r.text = text
        style(r, size, bold, color, face)
    if space is not None:
        p.line_spacing = space
    return tf


def bullets(sl, x, y, w, h, items, size=11.5, color=INK, gap=6, space=1.22):
    """items = [str | (str, bold_head)]; a leading bold lead-in before a space."""
    _, tf = textbox(sl, x, y, w, h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = space
        if i:
            p.space_before = Pt(gap)
        head, rest = (item if isinstance(item, tuple) else (None, item))
        if head:
            r = p.add_run(); r.text = head + "　"
            style(r, size, True, color)
        r = p.add_run(); r.text = rest
        style(r, size, False, color)
    return tf


def line(sl, x1, y1, x2, y2, color, width=1.0, dash=None):
    ln = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        el = ln.line._get_or_add_ln()
        el.append(el.makeelement(qn("a:prstDash"), {"val": dash}))
    return ln


def shape(sl, kind, x, y, w, h, fill=None, outline=None, ow=1.0):
    sh = sl.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if outline is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = outline
        sh.line.width = Pt(ow)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return sh


def box(sl, x, y, w, h, fill, text=None, size=11.5, bold=False,
        color=BLACK, face=LAT, outline=None, ow=1.0):
    sh = shape(sl, MSO_SHAPE.RECTANGLE, x, y, w, h, fill=fill, outline=outline, ow=ow)
    if text:
        r = sh.text_frame.paragraphs[0].add_run()
        r.text = text
        style(r, size, bold, color, face)
    return sh


# ------------------------------------------------------------ page frame ----
def page(prs, P, headline, eyebrow=None, deck=None, page_no=None, footer=None):
    """Standard content page: eyebrow, Arial Black headline, optional deck line.

    Returns (slide, y) where y is the top of the free content area.
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])       # blank
    if eyebrow:
        say(sl, L, 0.36, 10.0, 0.24, eyebrow, 11.5, True, P["key"])
    say(sl, L, 0.62, W, 0.64, headline, 24, True, BLACK, face=DISPLAY)
    y = 1.42
    if deck:
        say(sl, L, 1.34, W, 0.28, deck, 12.5, False, P["muted"])
        y = 1.80
    if footer:
        say(sl, L, FOOT_Y, 7.0, 0.22, footer, 9, False, P["muted"])
    if page_no is not None:
        say(sl, 12.30, FOOT_Y, 0.62, 0.22, str(page_no), 9, False, P["muted"],
            align=PP_ALIGN.RIGHT)
    return sl, y


def base_from_template(potx_path):
    """Strip every demo slide from an Accenture starter pack, keeping the master."""
    work = tempfile.mkdtemp(prefix="acnbase-")
    with zipfile.ZipFile(potx_path) as z:
        z.extractall(work)
    pres = os.path.join(work, "ppt", "presentation.xml")
    xml = open(pres, encoding="utf8").read()
    open(pres, "w", encoding="utf8").write(
        re.sub(r"<p:sldIdLst>.*?</p:sldIdLst>", "", xml, flags=re.S))
    ct = os.path.join(work, "[Content_Types].xml")
    types = open(ct, encoding="utf8").read()
    types = types.replace("presentationml.template.main+xml",
                          "presentationml.presentation.main+xml")
    types = re.sub(r'<Override PartName="/ppt/(slides|notesSlides)/[^"]*"[^>]*/>', "", types)
    open(ct, "w", encoding="utf8").write(types)
    for sub in ("slides", "notesSlides"):
        shutil.rmtree(os.path.join(work, "ppt", sub), ignore_errors=True)
    rels = os.path.join(work, "ppt", "_rels", "presentation.xml.rels")
    r = open(rels, encoding="utf8").read()
    open(rels, "w", encoding="utf8").write(
        re.sub(r'<Relationship[^>]*Target="slides/[^"]*"[^>]*/>', "", r))
    out = os.path.join(work, "_base.pptx")
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(work):
            for name in files:
                full = os.path.join(root, name)
                if full != out:
                    z.write(full, os.path.relpath(full, work))
    return out
