# -*- coding: utf-8 -*-
"""JIRA文書管理規程 — 説明用スライド（ACN風レイアウト）.

Usage:
    python build_deck.py --brand daicel [-o out.pptx]
    python build_deck.py --brand acn    [-o out.pptx]

規程本文（第1条〜第7条）を、結論を見出しに立てるコンサル型の構成に組み替えたもの。
版面は acn_kit の共通グリッドに従う。
"""
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from acn_kit import (BRANDS, BLACK, WHITE, INK, OK, WARN, DENY, OK_BG, WARN_BG,
                     DENY_BG, JP, LAT, DISPLAY, L, R, W, FOOT_Y,
                     say, rich, bullets, line, shape, box, page)

DOC = "JIRA文書管理規程"
FOOT = "生産準備部門　運用規程 ｜ JIRA文書管理規程　第1.0版"


# ---------------------------------------------------------------- 1. 表紙 ----
def cover(prs, P):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    box(sl, 0, 0, 13.333, 7.5, P["key"])                      # full bleed
    say(sl, 1.10, 2.05, 10.0, 0.30, "生産準備部門　運用規程", 13, True, WHITE)
    say(sl, 1.10, 2.52, 10.5, 0.95, DOC, 46, True, WHITE, face=DISPLAY)
    line(sl, 1.10, 3.72, 4.30, 3.72, WHITE, 2.0)
    say(sl, 1.10, 3.95, 8.6, 0.80,
        "プロジェクト管理に関する記録をタスク単位で保管し、\n"
        "社外との合意事項を後日確認できる状態を維持する。",
        14, False, WHITE, space=1.35)
    meta = [("制定日", "2026年9月4日", 2.30), ("版数", "第1.0版", 2.30),
            ("適用対象", "生準メンバ（今後、関連部署へ拡大）", 5.20)]
    x = 1.10
    for i, (k, v, cw) in enumerate(meta):
        if i:
            line(sl, x - 0.34, 5.72, x - 0.34, 6.24, WHITE, 0.75)
        say(sl, x, 5.72, cw - 0.40, 0.22, k, 10, True, WHITE)
        say(sl, x, 5.98, cw - 0.40, 0.26, v, 12, True, WHITE)
        x += cw
    say(sl, 1.10, 6.72, 8.0, 0.22,
        "※ 文書番号・管理部署は未記入（要確定）", 10, False, WHITE)


# ------------------------------------------------------------ 2. 全体像 ----
def overview(prs, P, n):
    sl, y = page(prs, P, "本規程は、記録の保管先と登録可否を定める7条で構成する",
                 eyebrow="全体像", deck="第1条〜第7条の関係。第3条・第6条が登録可否の判断軸、第5条が保管先の切り分けにあたる。",
                 page_no=n, footer=FOOT)
    items = [
        ("第1条", "目的", "個人のメールではなくチケットに残す"),
        ("第2条", "適用範囲・公開範囲", "生準メンバ → 関連部署へ段階的に公開"),
        ("第3条", "保管対象", "社外とのやり取りを最優先で登録する"),
        ("第4条", "社外コミュニケーション", "メールは原本のまま、口頭は確認メール"),
        ("第5条", "JIRA／Winchillの分担", "図面と版管理はWinchillに寄せる"),
        ("第6条", "情報区分と取扱い", "機密以上はリンクのみ、添付しない"),
        ("第7条", "登録前チェックリスト", "登録前に5点を確認する"),
    ]
    cw, gap = 2.98, 0.19
    for i, (no, title, sub) in enumerate(items):
        col, row = i % 4, i // 4
        x = L + col * (cw + gap)
        yy = y + row * 1.98
        box(sl, x, yy, cw, 1.78, P["panel"])
        box(sl, x, yy, 0.90, 0.32, P["key"], no, 10, True, WHITE)
        say(sl, x + 0.18, yy + 0.54, cw - 0.36, 0.52, title, 13.5, True, P["dark"], space=1.2)
        say(sl, x + 0.18, yy + 1.16, cw - 0.36, 0.52, sub, 10.5, False, P["muted"], space=1.28)
    box(sl, L, y + 4.06, W, 0.62, P["tint"])
    rich(sl, L + 0.22, y + 4.06, W - 0.44, 0.62,
         [("判断の軸　", 11.5, True, P["dark"]),
          ("「何を登録するか」＝第3条・第6条　／　「どこに保管するか」＝第5条　／　"
           "「どう残すか」＝第4条", 11.5, False, INK)],
         anchor=MSO_ANCHOR.MIDDLE)


# -------------------------------------------------------------- 3. 目的 ----
def purpose(prs, P, n):
    sl, y = page(prs, P, "記録は個人のメールボックスではなく、チケットに保管する",
                 eyebrow="第1条　目的", deck="個人の受信箱のみに存在する情報は、組織の記録として扱わない。",
                 page_no=n, footer=FOOT)
    goals = [("01", "認識の相違を防ぐ", "社外との協議・合意の経緯を確実に保全する。"),
             ("02", "属人化を避ける", "担当者の交代・不在時も、第三者が経緯を追跡できる。"),
             ("03", "事実を再構成する", "問題発生時に、事実関係を速やかに再構成できる。")]
    cw, gap = 4.03, 0.20
    for i, (no, t, d) in enumerate(goals):
        x = L + i * (cw + gap)
        box(sl, x, y, cw, 2.66, P["panel"])
        say(sl, x + 0.26, y + 0.30, 1.4, 0.58, no, 34, True, P["light"], face=DISPLAY)
        say(sl, x + 0.26, y + 1.06, cw - 0.52, 0.34, t, 16, True, P["dark"])
        say(sl, x + 0.26, y + 1.56, cw - 0.52, 0.80, d, 12, False, INK, space=1.35)
    ny = y + 3.02
    box(sl, L, ny, W, 1.32, P["tint"])
    box(sl, L, ny, 0.10, 1.32, P["key"])
    say(sl, L + 0.36, ny + 0.22, 2.0, 0.28, "原則", 14, True, P["dark"], face=DISPLAY)
    say(sl, L + 0.36, ny + 0.66, W - 0.72, 0.46,
        "記録は個人のメールボックスではなく、当該タスクのチケットに保管する。"
        "個人の受信箱のみに存在する情報は、組織の記録として扱わない。",
        13, False, INK, space=1.3)


# ------------------------------------------------------------ 4. 公開範囲 ----
def scope(prs, P, n):
    sl, y = page(prs, P, "公開範囲は生準メンバから関連部署へ段階的に広げる",
                 eyebrow="第2条　適用範囲および公開範囲",
                 deck="関連部署への公開を前提に、登録内容は当該範囲で閲覧されることを想定して記載する。",
                 page_no=n, footer=FOOT)
    stages = [("現行", "生準メンバのみ", "試行運用期間", P["out_fill"], P["out_fg"], P["muted"]),
              ("今後", "関連部署に公開する", "公開時期は別途通知する", P["key"], WHITE, WHITE)]
    aw, gap = 5.20, 1.10
    for i, (tag, who, note, fill, fg, sub) in enumerate(stages):
        x = L + 0.55 + i * (aw + gap)
        shape(sl, MSO_SHAPE.PENTAGON, x, y, aw, 1.50, fill=fill)
        say(sl, x + 0.42, y + 0.24, aw - 1.3, 0.26, tag, 11, True, fg)
        say(sl, x + 0.42, y + 0.58, aw - 1.3, 0.34, who, 18, True, fg, face=DISPLAY)
        say(sl, x + 0.42, y + 1.04, aw - 1.3, 0.24, note, 10.5, False, sub)
        if i == 0:
            tri = shape(sl, MSO_SHAPE.ISOSCELES_TRIANGLE, x + aw + 0.28, y + 0.56,
                        0.46, 0.40, fill=P["light"])
            tri.rotation = 90

    ly = y + 1.86
    say(sl, L, ly, 8.0, 0.24, "公開拡大に備えて、いまから徹底すること", 11.5, True, P["key"])
    prep = [("社外情報は原本で残す", "要約ではなく原本を保管し、閲覧者が経緯を追える状態にする。"),
            ("個人情報を含めない", "業務上不要な連絡先・私的な内容は登録前に取り除く。"),
            ("機密はリンクのみ", "機密以上はWinchill等に保管し、JIRAには参照先だけを書く。")]
    cw2, gap2 = 4.03, 0.20
    for i, (t, d) in enumerate(prep):
        x = L + i * (cw2 + gap2)
        box(sl, x, ly + 0.34, cw2, 1.34, P["panel"])
        box(sl, x, ly + 0.34, 0.06, 1.34, P["key"])
        say(sl, x + 0.24, ly + 0.54, cw2 - 0.48, 0.28, t, 13, True, P["dark"])
        say(sl, x + 0.24, ly + 0.92, cw2 - 0.48, 0.62, d, 10.5, False, INK, space=1.3)

    ny = ly + 2.00
    box(sl, L, ny, W, 0.84, P["tint"])
    rich(sl, L + 0.26, ny + 0.14, W - 0.52, 0.58,
         [("留意点　", 12, True, P["dark"]),
          ("関連部署に閲覧されて差し支えない表現か、登録の都度確認する"
           "（第7条のチェック項目に対応）。", 12, False, INK)], space=1.32)


# ------------------------------------------------------------ 5. 保管対象 ----
def targets(prs, P, n):
    sl, y = page(prs, P, "社外とのやり取りを最優先で登録し、図面と個人情報は登録しない",
                 eyebrow="第3条　保管対象",
                 deck="プロジェクト管理に関連する情報を、該当するチケットに保管する。",
                 page_no=n, footer=FOOT)
    cw = 6.15
    # 登録するもの
    box(sl, L, y, cw, 4.56, P["panel"])
    box(sl, L, y, cw, 0.42, P["key"])
    say(sl, L + 0.22, y + 0.09, cw - 0.44, 0.26, "登録するもの", 12.5, True, WHITE, face=DISPLAY)
    rows = [("社外とのやり取り", "顧客・サプライヤとの依頼、回答、条件調整、合意", True),
            ("決定事項", "方針決定および重要な決議内容", False),
            ("進捗報告", "定期報告、マイルストーンの達成状況", False),
            ("要件確認", "仕様決定、変更依頼の承認", False),
            ("議事録", "打合せの記録および議論の経過", False),
            ("ステータス", "課題・リスク・対応状況", False)]
    ry = y + 0.72
    for label, desc, top in rows:
        say(sl, L + 0.24, ry, 2.15, 0.26, label, 12, True,
            P["key"] if top else P["dark"])
        say(sl, L + 2.48, ry + 0.02, cw - 2.72, 0.46, desc, 10.5, False, INK, space=1.28)
        if top:
            say(sl, L + 0.24, ry + 0.28, 2.15, 0.20, "最重要", 9, True, P["key"])
            ry += 0.22
        ry += 0.60
    # 登録しないもの
    x2 = L + cw + 0.20
    box(sl, x2, y, cw, 4.56, DENY_BG)
    box(sl, x2, y, cw, 0.42, DENY)
    say(sl, x2 + 0.22, y + 0.09, cw - 0.44, 0.26, "登録しないもの", 12.5, True, WHITE, face=DISPLAY)
    nos = [("図面・設計図・CADファイル", "Winchillにて保管する（第5条）"),
           ("業務上不要な個人情報", "個人携帯番号、自宅住所、私的な内容")]
    ry = y + 0.78
    for label, desc in nos:
        say(sl, x2 + 0.24, ry, cw - 0.48, 0.28, label, 12.5, True, DENY)
        say(sl, x2 + 0.24, ry + 0.34, cw - 0.48, 0.26, desc, 10.5, False, INK)
        ry += 1.00
    box(sl, x2 + 0.24, y + 2.94, cw - 0.48, 1.42, WHITE)
    say(sl, x2 + 0.46, y + 3.12, cw - 0.92, 0.32, "判断に迷ったら", 11.5, True, P["dark"])
    say(sl, x2 + 0.46, y + 3.52, cw - 0.92, 0.72,
        "登録前に第7条のチェックリストで確認する。"
        "機密区分に該当するものは、JIRAには参照リンクのみを記載する。",
        10.5, False, INK, space=1.32)


# -------------------------------------------------- 6. 社外コミュニケーション ----
def external(prs, P, n):
    sl, y = page(prs, P, "社外メールは原本のまま保管し、口頭決定は確認メールで残す",
                 eyebrow="第4条　社外コミュニケーションの記録",
                 deck="要約のみの記載は認めない。依頼と回答は対で保管する。",
                 page_no=n, footer=FOOT)
    rules = [("01", "原本のまま保管する",
              "社外とのメールは、該当チケットに原本のまま保管する。要約のみの記載は認めない。"),
             ("02", "依頼と回答を対で残す",
              "依頼と回答は対で保管し、合意内容を読み取れる状態とする。"),
             ("03", "訂正は追記で残す",
              "訂正が生じた場合、元の記録は削除せず、コメントにより追記する。")]
    ry = y
    for no, t, d in rules:
        box(sl, L, ry, W, 1.05, P["panel"])
        box(sl, L, ry, 1.00, 1.05, P["key"], no, 20, True, WHITE, face=DISPLAY)
        say(sl, L + 1.26, ry + 0.36, 3.60, 0.32, t, 14, True, P["dark"])
        say(sl, L + 5.00, ry + 0.30, W - 5.30, 0.46, d, 11.5, False, INK, space=1.3)
        ry += 1.18
    ny = ry + 0.24
    box(sl, L, ny, W, 1.10, WARN_BG)
    box(sl, L, ny, 0.10, 1.10, WARN)
    say(sl, L + 0.36, ny + 0.18, 5.0, 0.28, "電話・口頭で決定した場合", 13.5, True, WARN, face=DISPLAY)
    say(sl, L + 0.36, ny + 0.60, W - 0.72, 0.44,
        "内容の確認メールを相手方へ送付し、当該メールをチケットに保管する。"
        "本項の徹底が、認識相違に対する最も有効な予防措置となる。",
        12.5, False, INK, space=1.3)


# ------------------------------------------------------ 7. JIRA / Winchill ----
def systems(prs, P, n):
    sl, y = page(prs, P, "JIRAはコミュニケーション記録、Winchillは図面と版管理を担う",
                 eyebrow="第5条　JIRAおよびWinchillの役割分担",
                 deck="JIRAでは版の履歴を都度確認できないため、バージョン管理はWinchillに一元化する。",
                 page_no=n, footer=FOOT)
    cols = [("JIRA", P["key"], WHITE,
             ["決定事項", "進捗報告", "議事録", "ステータス", "要件確認"],
             "プロジェクト管理／コミュニケーション記録",
             "生準メンバ（今後、関連部署）",
             "版の履歴を都度確認できない。最新版の所在は、"
             "参照リンクまたはWinchill No.で示す。"),
            ("Winchill", P["dark"], WHITE,
             ["図面・設計図", "CADファイル", "技術仕様書", "製造情報", "その他機密資料"],
             "機密情報管理／設計・技術情報の版管理",
             "権限者のみ（別途制御）",
             "版管理はここに一元化する。JIRA側からは"
             "参照リンクで辿れる状態にしておく。")]
    cw = 6.15
    for i, (name, fill, fg, keeps, use, acc, note) in enumerate(cols):
        x = L + i * (cw + 0.20)
        box(sl, x, y, cw, 3.85, P["panel"])
        box(sl, x, y, cw, 0.52, fill)
        say(sl, x + 0.24, y + 0.13, cw - 0.48, 0.30, name, 16, True, fg, face=DISPLAY)
        say(sl, x + 0.24, y + 0.70, cw - 0.48, 0.20, "保管対象", 10, True, P["muted"])
        cy = y + 0.98
        for k in keeps:
            box(sl, x + 0.24, cy, 2.05, 0.30, fill, k, 10.5, True, fg)
            cy += 0.38
        say(sl, x + 2.55, y + 0.98, cw - 2.80, 0.20, "用途", 10, True, P["muted"])
        say(sl, x + 2.55, y + 1.22, cw - 2.80, 0.46, use, 11, False, INK, space=1.3)
        say(sl, x + 2.55, y + 1.92, cw - 2.80, 0.20, "アクセス", 10, True, P["muted"])
        say(sl, x + 2.55, y + 2.16, cw - 2.80, 0.46, acc, 11, False, INK, space=1.3)
        line(sl, x + 0.24, y + 2.86, x + cw - 0.24, y + 2.86, P["rule"], 0.75)
        say(sl, x + 0.24, y + 3.00, cw - 0.48, 0.20, "運用上の注意", 10, True, P["muted"])
        say(sl, x + 0.24, y + 3.24, cw - 0.48, 0.52, note, 11, False, INK, space=1.32)
    ny = y + 4.05
    box(sl, L, ny, W, 0.90, P["tint"])
    rich(sl, L + 0.26, ny + 0.15, W - 0.52, 0.62,
         [("図面・CADファイルの取扱い　", 12, True, P["dark"]),
          ("JIRAには登録せず、Winchillにて保管する。"
           "JIRAには参照リンクまたはWinchill No.のみを記載する。", 12, False, INK)], space=1.32)


# ---------------------------------------------------------- 8. 情報区分 ----
def classes(prs, P, n):
    sl, y = page(prs, P, "機密以上はJIRAに登録せず、参照リンクのみを記載する",
                 eyebrow="第6条　情報区分と取扱い",
                 deck="機密区分の情報は添付ファイルとしては登録しない。",
                 page_no=n, footer=FOOT)
    rows = [("一般", "Public", "可", OK, OK_BG, "JIRA"),
            ("社内", "Internal", "可", OK, OK_BG, "JIRA（アクセス権限を設定する）"),
            ("機密", "Confidential", "不可", DENY, DENY_BG, "Winchill／専用ストレージ　※参照リンクのみ"),
            ("極秘", "Secret", "不可", DENY, DENY_BG, "限定者のみアクセス可能な場所")]
    hy = y
    say(sl, L + 0.20, hy, 3.10, 0.24, "情報区分", 10.5, True, P["muted"])
    say(sl, L + 3.30, hy, 1.10, 0.24, "JIRAへの登録", 10.5, True, P["muted"],
        align=PP_ALIGN.CENTER)
    say(sl, L + 5.60, hy, W - 5.80, 0.24, "保管先", 10.5, True, P["muted"])
    ry = hy + 0.34
    for jp, en, ok, col, bg, dest in rows:
        box(sl, L, ry, W, 0.80, bg)
        say(sl, L + 0.20, ry + 0.14, 2.90, 0.28, jp, 15, True, col, face=DISPLAY)
        say(sl, L + 0.20, ry + 0.47, 2.90, 0.20, en, 9.5, False, P["muted"])
        box(sl, L + 3.30, ry + 0.22, 1.10, 0.36, col, ok, 12, True, WHITE, face=DISPLAY)
        say(sl, L + 5.60, ry + 0.28, W - 5.80, 0.28, dest, 12, False, INK)
        ry += 0.88
    ny = ry + 0.14
    box(sl, L, ny, W, 0.78, P["panel"])
    rich(sl, L + 0.26, ny + 0.12, W - 0.52, 0.54,
         [("運用　", 12, True, P["dark"]),
          ("機密区分の情報を参照する必要がある場合は、JIRAには外部リンクのみを記載する。"
           "添付ファイルとしては登録しない。", 12, False, INK)], space=1.32)


# -------------------------------------------------------- 9. チェックリスト ----
def checklist(prs, P, n):
    sl, y = page(prs, P, "登録前に5点を確認する",
                 eyebrow="第7条　登録前チェックリスト",
                 deck="JIRAに文書を登録する前に、以下を確認する。",
                 page_no=n, footer=FOOT)
    items = [("図面・CADファイルではないか", "Winchillで保管すべきものではないか"),
             ("顧客の個人情報・連絡先が含まれていないか", "個人携帯番号・自宅住所・私的な内容"),
             ("必要なアクセス権限が設定されているか", "社内区分は権限設定のうえ登録する"),
             ("ファイルサイズは50MB以下か", "超過する場合は分割または別保管を検討する"),
             ("関連部署に閲覧されても差し支えない内容か", "公開範囲の拡大を前提に判断する")]
    ry = y
    for i, (q, note) in enumerate(items, 1):
        box(sl, L, ry, W, 0.72, P["panel"] if i % 2 else WHITE)
        shape(sl, MSO_SHAPE.RECTANGLE, L + 0.28, ry + 0.22, 0.28, 0.28,
              fill=WHITE, outline=P["key"], ow=1.25)
        say(sl, L + 0.86, ry + 0.10, 0.50, 0.24, f"0{i}", 10, True, P["light"], face=DISPLAY)
        say(sl, L + 0.86, ry + 0.32, 6.40, 0.28, q, 13, True, P["dark"])
        say(sl, L + 7.50, ry + 0.24, W - 7.70, 0.28, note, 11, False, P["muted"])
        ry += 0.80
    ny = ry + 0.10
    box(sl, L, ny, W, 0.72, WARN_BG)
    box(sl, L, ny, 0.10, 0.72, WARN)
    rich(sl, L + 0.34, ny + 0.12, W - 0.70, 0.50,
         [("図面・設計図の場合　", 12, True, WARN),
          ("Winchillに登録のうえ、JIRAには参照リンクまたはWinchill No.を記載する。",
           12, False, INK)], space=1.3)


# ----------------------------------------------------------- 10. まとめ ----
def closing(prs, P, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    box(sl, 0, 0, 13.333, 7.5, P["key"])
    say(sl, L + 0.68, 0.80, 10.0, 0.26, "まとめ", 12, True, WHITE)
    say(sl, L + 0.68, 1.20, 11.4, 1.10,
        "記録の所在を一つに定めることが、\n認識の相違と属人化を防ぐ",
        34, True, WHITE, face=DISPLAY, space=1.22)
    line(sl, L + 0.68, 2.70, L + 3.60, 2.70, WHITE, 2.0)
    pts = [("残す場所", "個人のメールボックスではなく、当該タスクのチケットに保管する。"),
           ("残し方", "社外メールは原本のまま。口頭決定は確認メールを送り、それを保管する。"),
           ("分けるもの", "図面・CADと機密以上はWinchillへ。JIRAには参照リンクのみを記載する。")]
    cw, gap = 3.72, 0.30
    for i, (t, d) in enumerate(pts):
        x = L + 0.68 + i * (cw + gap)
        line(sl, x, 3.20, x + cw - 0.30, 3.20, WHITE, 0.75)
        say(sl, x, 3.36, cw - 0.30, 0.30, t, 15, True, WHITE, face=DISPLAY)
        say(sl, x, 3.78, cw - 0.30, 0.90, d, 11.5, False, WHITE, space=1.35)
    box(sl, L + 0.68, 5.10, 11.4, 0.86, WHITE)
    rich(sl, L + 0.92, 5.24, 11.0, 0.60,
         [("未決事項　", 12, True, P["dark"]),
          ("① 文書番号の採番　　② 管理部署の確定　　③ 関連部署への公開時期（別途通知）",
           12, False, INK)], space=1.3)
    say(sl, L + 0.68, 6.44, 11.4, 0.24,
        "附則　本規程は制定日より実施する。改定は管理部署の承認を経て行う。",
        11, False, WHITE)
    say(sl, 12.30, FOOT_Y, 0.62, 0.22, str(n), 9, False, WHITE, align=PP_ALIGN.RIGHT)


def build(brand, output):
    P = BRANDS[brand]
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    cover(prs, P)
    for i, fn in enumerate([overview, purpose, scope, targets, external,
                            systems, classes, checklist], start=2):
        fn(prs, P, i)
    closing(prs, P, 10)
    prs.save(output)
    print("saved", output, "-", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--brand", choices=sorted(BRANDS), default="daicel")
    ap.add_argument("-o", "--output")
    a = ap.parse_args()
    build(a.brand, a.output or f"jira_document_policy_{a.brand}.pptx")
