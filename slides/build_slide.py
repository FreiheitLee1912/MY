# -*- coding: utf-8 -*-
"""Build the "JIRA登録タイミング" slide in either Accenture or DAICEL colours.

Usage:
    python build_slide.py --brand daicel [-o out.pptx]
    python build_slide.py --brand acn --template <Acc_PPT_IMP_StarterPack....potx>

Both brands share one layout system, lifted from the Accenture 2020 starter pack:
0.42" side margins over a 12.50" text column, an Arial Black headline set flush
top-left with no rule under it, and a conclusion ("So What") band closing the page.
Only the palette and the page furniture change between brands.

  acn     - built on the starter pack's own "Long Headline Only" layout, so the
            master furniture (the purple ">" mark, the footer band) is inherited.
  daicel  - built on a bare page in DAICEL blue. The Accenture master is dropped
            deliberately: its ">" mark is Accenture's, and purple would fight the
            blue. Drop your own logo into the footer.

Requires: python-pptx
"""
import argparse
import os
import re
import shutil
import tempfile
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def C(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


BLACK, WHITE = C("000000"), C("FFFFFF")
INK = C("333333")

BRANDS = {
    # Accenture 2020: accent1 A100FF, accent2 7500C0, accent3 460073,
    # accent5 BE82FF, dk2 96968C, lt2 E6E6DC
    "acn": dict(
        key=C("A100FF"), mid=C("7500C0"), dark=C("460073"), light=C("BE82FF"),
        tint=C("F4E9FF"), muted=C("70706A"),
        out_fill=C("E6E6DC"), out_fg=C("70706A"),
        pre_fill=C("F4E9FF"), pre_fg=C("460073"), pre_line=C("A100FF"),
        post_fills=(C("A100FF"), C("7500C0")), post_fg=WHITE,
        after_sop=C("460073"),                 # "PCR都度起票"
        open_fill=C("E6E6DC"), open_fg=C("460073"),
        sowhat_fill=C("A100FF"), sowhat_fg=WHITE,
        axis=C("96968C"), track=C("96968C"),
    ),
    # DAICEL: corporate blue, with the source slide's green kept for the
    # post-SOP phases so the original colour semantics survive the reskin.
    "daicel": dict(
        key=C("0068B7"), mid=C("1B6EC2"), dark=C("1F4E79"), light=C("9DC3E6"),
        tint=C("E8F1FA"), muted=C("6B7280"),
        out_fill=C("E7E7E7"), out_fg=C("767676"),
        pre_fill=C("E8F1FA"), pre_fg=C("1F4E79"), pre_line=C("0068B7"),
        post_fills=(C("4C9A2A"), C("3D7C21")), post_fg=WHITE,
        after_sop=C("3D7C21"),
        open_fill=C("FBF4E0"), open_fg=C("0068B7"),
        sowhat_fill=C("0068B7"), sowhat_fg=WHITE,
        axis=C("A6A6A6"), track=C("A6A6A6"),
    ),
}

JP, LAT, DISPLAY = "Meiryo", "Arial", "Arial Black"

# ---- ACN grid ---------------------------------------------------------------
L, R = 0.42, 12.92
TRACK_L = 1.62
N = 9
PITCH = (R - TRACK_L) / N
CHEV_W = PITCH - 0.045
bx = lambda i: TRACK_L + PITCH * i
cx = lambda i: TRACK_L + PITCH * i + PITCH / 2
SOP_X, INF_X, INIT_X = bx(6), cx(2), cx(3)
AXIS_Y = 3.92

PHASES = ["引合", "提案", "仕様検討", "投資", "量産準備", "検証", "量産", "補給品", "EOP"]


def base_from_template(potx_path):
    """Strip every demo slide from the starter pack, keeping master + layouts."""
    work = tempfile.mkdtemp(prefix="acnbase-")
    with zipfile.ZipFile(potx_path) as z:
        z.extractall(work)
    pres = os.path.join(work, "ppt", "presentation.xml")
    xml = open(pres, encoding="utf8").read()
    open(pres, "w", encoding="utf8").write(
        re.sub(r"<p:sldIdLst>.*?</p:sldIdLst>", "", xml, flags=re.S))
    ct = os.path.join(work, "[Content_Types].xml")
    types = open(ct, encoding="utf8").read()
    types = types.replace("presentationml.template.main+xml",
                          "presentationml.presentation.main+xml")
    types = re.sub(r'<Override PartName="/ppt/(slides|notesSlides)/[^"]*"[^>]*/>', "", types)
    open(ct, "w", encoding="utf8").write(types)
    for sub in ("slides", "notesSlides"):
        shutil.rmtree(os.path.join(work, "ppt", sub), ignore_errors=True)
    rels = os.path.join(work, "ppt", "_rels", "presentation.xml.rels")
    r = open(rels, encoding="utf8").read()
    open(rels, "w", encoding="utf8").write(
        re.sub(r'<Relationship[^>]*Target="slides/[^"]*"[^>]*/>', "", r))
    out = os.path.join(work, "_base.pptx")
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(work):
            for name in files:
                full = os.path.join(root, name)
                if full != out:
                    z.write(full, os.path.relpath(full, work))
    return out


# ---- drawing helpers --------------------------------------------------------
def style(run, size, bold=False, color=BLACK, face=LAT, ea=JP):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, face
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag, val in (("a:ea", ea), ("a:cs", face)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", val)


def textbox(sl, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tb, tf


def say(sl, x, y, w, h, text, size, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, face=LAT, anchor=MSO_ANCHOR.TOP):
    _, tf = textbox(sl, x, y, w, h, align, anchor)
    r = tf.paragraphs[0].add_run()
    r.text = text
    style(r, size, bold, color, face)
    return tf


def line(sl, x1, y1, x2, y2, color, width=1.0, dash=None):
    ln = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        el = ln.line._get_or_add_ln()
        el.append(el.makeelement(qn("a:prstDash"), {"val": dash}))
    return ln


def shape(sl, kind, x, y, w, h, fill=None, outline=None, ow=1.0):
    sh = sl.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if outline is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = outline
        sh.line.width = Pt(ow)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return sh


# ---- the slide --------------------------------------------------------------
def build(brand, template, output):
    P = BRANDS[brand]
    KEY, MID, DARK, LIGHT = P["key"], P["mid"], P["dark"], P["light"]
    MUTED, TINT = P["muted"], P["tint"]

    if brand == "acn":
        if not template:
            raise SystemExit("--brand acn needs --template <starter pack .potx>")
        prs = Presentation(base_from_template(template))
        layout = next(l for l in prs.slide_master.slide_layouts
                      if l.name == "Long Headline Only")
        sl = prs.slides.add_slide(layout)
        head = sl.shapes.title
        head.left, head.top = Inches(L), Inches(0.62)
        head.width, head.height = Inches(12.50), Inches(0.66)
        tf = head.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        r = tf.paragraphs[0].add_run()
        r.text = "INFは開発企画書発行時、イニシは仕様FIX時にJIRAへ登録する"
        style(r, 27, True, BLACK, face=DISPLAY)
    else:
        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        sl = prs.slides.add_slide(prs.slide_layouts[6])   # blank
        say(sl, L, 0.62, 12.50, 0.66,
            "INFは開発企画書発行時、イニシは仕様FIX時にJIRAへ登録する",
            27, True, BLACK, face=DISPLAY)

    # eyebrow
    say(sl, L, 0.36, 9.0, 0.24,
        "JIRA登録タイミング ｜ INF・イニシ別の登録トリガー", 11.5, True, KEY)
    # deck
    say(sl, L, 1.34, 12.0, 0.28,
        "登録判断の位置を入口・出口に分け、SOPを境界に管理区間を切り替える", 13.5, False, MUTED)

    # --- definition band ---
    shape(sl, MSO_SHAPE.RECTANGLE, L, 1.78, R - L, 0.70, fill=TINT)
    _, tf = textbox(sl, L + 0.22, 1.90, 11.9, 0.24)
    p = tf.paragraphs[0]
    for t, sz, b, c in [("登録判断", 12, True, DARK), ("     ", 12, False, BLACK),
                        ("INF ＝ 開発企画書 発行時", 12.5, True, BLACK),
                        ("     ｜     ", 12.5, False, P["axis"]),
                        ("イニシ ＝ 仕様FIX時", 12.5, True, BLACK),
                        ("     ｜     ", 12.5, False, P["axis"]),
                        ("登録以降はJIRAで進捗・課題を管理", 12.5, True, MID)]:
        r = p.add_run(); r.text = t; style(r, sz, b, c)
    say(sl, L + 0.22, 2.18, 11.9, 0.22,
        "登録前は個別に管理し、登録後は案件の進捗・課題をJIRA上で追跡する。", 10.5, False, MUTED)

    # --- lifecycle ---
    say(sl, L, 2.64, 4.0, 0.24, "プロジェクトライフサイクル", 11.5, True, KEY)
    for txt, x0, x1, col in [("SOP前", bx(4), bx(6), MID),
                             ("SOP", SOP_X - 0.30, SOP_X + 0.30, DARK),
                             ("SOP後", bx(6), bx(8), P["after_sop"])]:
        say(sl, x0, 2.64, x1 - x0, 0.24, txt, 10.5, True, col, align=PP_ALIGN.CENTER)
    line(sl, SOP_X, 2.92, SOP_X, 5.44, DARK, 1.25, dash="dash")

    for i, name in enumerate(PHASES):
        if i in (0, 1, 8):
            fill, fg, out = P["out_fill"], P["out_fg"], None
        elif i in (6, 7):
            fill, fg, out = P["post_fills"][i - 6], P["post_fg"], None
        else:
            fill, fg, out = P["pre_fill"], P["pre_fg"], P["pre_line"]
        ch = shape(sl, MSO_SHAPE.CHEVRON, bx(i), 3.00, CHEV_W, 0.44,
                   fill=fill, outline=out, ow=1.0)
        r = ch.text_frame.paragraphs[0].add_run(); r.text = name
        style(r, 11, True, fg)

    # --- milestones ---
    for x, label in [(INF_X, "開発企画書 発行"), (INIT_X, "仕様FIX")]:
        say(sl, x - 1.05, 3.56, 2.10, 0.22, label, 10.5, True, DARK, align=PP_ALIGN.CENTER)
    line(sl, TRACK_L, AXIS_Y, R, AXIS_Y, P["axis"], 0.75)
    for x in (INF_X, INIT_X):
        shape(sl, MSO_SHAPE.DIAMOND, x - 0.075, AXIS_Y - 0.075, 0.15, 0.15, fill=KEY)

    # --- swim lanes ---
    for name, y, trig in [("INF", 4.64, INF_X), ("イニシ", 5.26, INIT_X)]:
        say(sl, L, y - 0.13, 1.05, 0.26, name, 14, True, BLACK,
            align=PP_ALIGN.RIGHT, face=DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
        line(sl, TRACK_L, y, trig, y, P["track"], 1.25)
        line(sl, trig, y, R, y, KEY, 1.75)
        line(sl, trig, AXIS_Y, trig, y, LIGHT, 0.75, dash="sysDot")
        shape(sl, MSO_SHAPE.DIAMOND, trig - 0.075, y - 0.075, 0.15, 0.15, fill=KEY)
        for txt, x0, x1, sz, b, col in [
                ("個別管理", TRACK_L + 0.06, trig, 10.5, False, MUTED),
                ("JIRAで管理", trig, SOP_X, 11.5, True, MID),
                ("PCR都度起票", SOP_X, R, 11.5, True, P["after_sop"])]:
            say(sl, x0, y - 0.34, x1 - x0, 0.22, txt, sz, b, col, align=PP_ALIGN.CENTER)

    # --- open items ---
    shape(sl, MSO_SHAPE.RECTANGLE, L, 5.54, R - L, 0.60, fill=P["open_fill"])
    say(sl, L + 0.22, 5.65, 1.4, 0.22, "未決事項", 11.5, True, P["open_fg"])
    for txt, x0 in [("① 登録判断は全件か選別か（タイミング定義と判定基準の要否）", 1.90),
                    ("② 補給品の管理媒体（JIRA／Excel）と既存方針の整合", 7.55)]:
        say(sl, x0, 5.65, 5.35, 0.40, txt, 11, False, INK)

    # --- So What ---
    shape(sl, MSO_SHAPE.RECTANGLE, L, 6.32, R - L, 0.62, fill=P["sowhat_fill"])
    say(sl, L + 0.30, 6.44, 1.6, 0.26, "So What", 13, True, P["sowhat_fg"],
        face=DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
    say(sl, L + 2.05, 6.44, 10.3, 0.26,
        "INF／イニシを登録起点として定義し、SOP前後の管理責任を明確にする",
        13.5, True, P["sowhat_fg"], anchor=MSO_ANCHOR.MIDDLE)

    # --- footer ---
    foot_x = L + 0.32 if brand == "acn" else L      # ACN master draws ">" at 0.42
    say(sl, foot_x, 7.10, 6.0, 0.22, "JIRA登録タイミング ｜ トリガー定義", 9, False, MUTED)
    say(sl, 12.30, 7.10, 0.62, 0.22, "8", 9, False, MUTED, align=PP_ALIGN.RIGHT)

    sl.notes_slide.notes_text_frame.text = (
        "INFは開発企画書の発行時点、イニシは仕様FIX時点をJIRA登録のトリガーとする。"
        "登録前は個別管理、登録後はJIRA上で進捗・課題を追跡し、SOPを境界に量産・補給品は"
        "PCRを都度起票する運用へ切り替える。未決は登録判断の全件／選別の別と、"
        "補給品の管理媒体（JIRA／Excel）の整合。")

    prs.save(output)
    print("saved", output)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--brand", choices=sorted(BRANDS), default="daicel")
    ap.add_argument("--template", help="Accenture starter pack .potx (--brand acn only)")
    ap.add_argument("-o", "--output")
    a = ap.parse_args()
    build(a.brand, a.template, a.output or f"jira_registration_timing_{a.brand}.pptx")
