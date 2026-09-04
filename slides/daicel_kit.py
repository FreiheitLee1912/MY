# -*- coding: utf-8 -*-
"""Page system for the DAICEL standard template (2026_Standard_Template_WideScreen).

Measured from the template's own master and layouts:
  タイトル       x=0.92 y=0.19 w=10.30 h=0.60   (master default 40pt Meiryo UI;
                                                 overridden to 22pt for Japanese)
  本文領域       x=0.92 y=1.05 w=11.50 h=5.70
  日付           x=0.92 y=6.90 w=2.93 h=0.40
  注記(著作権)   x=4.07 y=6.90 w=5.19 h=0.40    8pt / 808080
  ページ番号     x=9.48 y=6.90 w=2.93 h=0.40
Master furniture (logo, top rules, bottom copyright band) is inherited, not redrawn.
Body face is Meiryo (the master defaults to Meiryo UI; this deck overrides it).
"""
import os
import re
import shutil
import tempfile
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import titus_kit
from titus_kit import (C, tint, BLACK, WHITE, style, textbox, say, rich, line,
                       cell)

# The DAICEL template's master specifies Meiryo UI; this deck uses Meiryo.
# Meiryo carries more line gap than Meiryo UI, so fixed-height cells here are
# sized with that extra leading in mind.
FONT = "Meiryo"
titus_kit.FONT = FONT

# ---- palette: DAICEL brand blue sampled from the template's own logo ---------
BLUE = C("0096D8")        # corporate blue
DARK = C("00558C")        # deep blue for emphasis
PALE = C("E5F4FC")        # label cells
GREY = C("E7E6E6")        # lt2
RULE = C("BFBFBF")        # table hairline - grey, not black (DAICEL is cleaner)
MUTED = C("595959")
FAINT = C("808080")
RED = C("C00000")
GREEN = C("00843D")
TAG_DECIDED = tint("ED7D31")   # accent2
TAG_OPEN = tint("7030A0")
NOTE_FILL = C("EAF4FB")
WARN_FILL = tint("ED7D31")

# ---- grid -------------------------------------------------------------------
L, R = 0.92, 12.42
W = R - L
BODY_TOP = 1.78
FOOT_Y = 6.90
DISCLAIMER = ("The entire contents of this publication are copyrighted by Daicel. "
              "It is forbidden to duplicate or alter this document, or to use its "
              "content for another purpose, without the express permission of Daicel.")


def base_from_template(path):
    """Strip the template's sample slides, keeping master, layouts and media."""
    work = tempfile.mkdtemp(prefix="dcbase-")
    with zipfile.ZipFile(path) as z:
        z.extractall(work)
    pres = os.path.join(work, "ppt", "presentation.xml")
    xml = open(pres, encoding="utf8").read()
    open(pres, "w", encoding="utf8").write(
        re.sub(r"<p:sldIdLst>.*?</p:sldIdLst>", "", xml, flags=re.S))
    ct = os.path.join(work, "[Content_Types].xml")
    types = open(ct, encoding="utf8").read()
    types = types.replace("presentationml.template.main+xml",
                          "presentationml.presentation.main+xml")
    types = re.sub(r'<Override PartName="/ppt/(slides|notesSlides)/[^"]*"[^>]*/>', "", types)
    open(ct, "w", encoding="utf8").write(types)
    for sub in ("slides", "notesSlides"):
        shutil.rmtree(os.path.join(work, "ppt", sub), ignore_errors=True)
    rels = os.path.join(work, "ppt", "_rels", "presentation.xml.rels")
    r = open(rels, encoding="utf8").read()
    open(rels, "w", encoding="utf8").write(
        re.sub(r'<Relationship[^>]*Target="slides/[^"]*"[^>]*/>', "", r))
    out = os.path.join(work, "_base.pptx")
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(work):
            for name in files:
                full = os.path.join(root, name)
                if full != out:
                    z.write(full, os.path.relpath(full, work))
    return out


def drop_placeholder(sl, idx):
    """Remove a cloned layout placeholder we are not going to fill."""
    for ph in list(sl.placeholders):
        if ph.placeholder_format.idx == idx:
            ph._element.getparent().remove(ph._element)


def secbar(sl, x, y, w, text, h=0.38, size=14, fill=None, fg=WHITE):
    """DAICEL section bar: brand blue, white label, no border."""
    return cell(sl, x, y, w, h, text, size, True, fg, fill or BLUE, None,
                align=PP_ALIGN.LEFT, pad=0.16)


def note(sl, x, y, w, h, label, body, fill=NOTE_FILL, size=12):
    cell(sl, x, y, w, h, "", fill=fill, border=None)
    parts = ([(label + "　", size, True, DARK)] if label else []) + \
            [(body, size, False, BLACK)]
    rich(sl, x + 0.20, y + 0.11, w - 0.40, h - 0.22, parts,
         anchor=MSO_ANCHOR.MIDDLE, space=1.32)


def foot(sl, page_no, date="2026/09/04"):
    say(sl, L, FOOT_Y + 0.06, 2.93, 0.28, date, 12, False, MUTED)
    say(sl, 4.07, FOOT_Y - 0.02, 5.19, 0.38, DISCLAIMER, 7, False, FAINT, space=1.12)
    say(sl, 9.48, FOOT_Y + 0.04, 2.93, 0.28, str(page_no), 14, False, MUTED,
        align=PP_ALIGN.RIGHT)


def page(prs, title, desc=None, tags=(), page_no=None, layout=1):
    """Content page on the template's タイトルとコンテンツ layout."""
    sl = prs.slides.add_slide(prs.slide_master.slide_layouts[layout])
    drop_placeholder(sl, 1)                     # the 28pt bullet body - we draw our own
    t = sl.shapes.title
    t.left, t.top = Inches(L), Inches(0.19)
    t.width, t.height = Inches(10.30), Inches(0.60)
    tf = t.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    style(r, 22, True, BLACK)
    if desc:
        say(sl, L, 1.02, 9.30, 0.56, desc, 12.5, False, MUTED, space=1.3)
    for i, (label, fill) in enumerate(reversed(list(tags))):
        cell(sl, R - 1.32 - i * 1.38, 1.02, 1.32, 0.30, label, 11, False, BLACK,
             fill, RULE, 0.75)
    foot(sl, page_no)
    return sl
