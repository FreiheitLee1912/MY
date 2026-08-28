# -*- coding: utf-8 -*-
"""Rebuild the "JIRA登録タイミング" slide in Accenture (ACN) house style.

Usage:
    python build_acn_slide.py <Accenture_StarterPack.potx> [out.pptx]

The design system is read from the Accenture 2020 starter-pack theme:
  purple A100FF / 7500C0 / 460073 / BE82FF / DCAFFF, warm grey 96968C / E6E6DC,
  Arial Black headlines over Arial body copy, 0.42" side margins, 12.50" column.
The slide is built on the pack's own "Long Headline Only" layout, so the master
furniture (the purple ">" mark, the footer band) is inherited, not redrawn.

Requires: python-pptx
"""
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def base_from_template(potx_path):
    """Strip every demo slide from the starter pack, keeping master + layouts."""
    work = tempfile.mkdtemp(prefix="acnbase-")
    with zipfile.ZipFile(potx_path) as z:
        z.extractall(work)
    pres = os.path.join(work, "ppt", "presentation.xml")
    with open(pres, encoding="utf8") as fh:
        xml = fh.read()
    xml = re.sub(r"<p:sldIdLst>.*?</p:sldIdLst>", "", xml, flags=re.S)
    with open(pres, "w", encoding="utf8") as fh:
        fh.write(xml)
    # a .potx declares a template content type; the output is a presentation
    ct = os.path.join(work, "[Content_Types].xml")
    with open(ct, encoding="utf8") as fh:
        types = fh.read()
    with open(ct, "w", encoding="utf8") as fh:
        fh.write(types.replace("presentationml.template.main+xml",
                               "presentationml.presentation.main+xml"))
    # drop slide-only parts the emptied deck no longer references
    for sub in ("ppt/slides", "ppt/notesSlides"):
        shutil.rmtree(os.path.join(work, *sub.split("/")), ignore_errors=True)
    for rels in ("ppt/_rels/presentation.xml.rels",):
        path = os.path.join(work, *rels.split("/"))
        with open(path, encoding="utf8") as fh:
            r = fh.read()
        r = re.sub(r'<Relationship[^>]*Target="slides/[^"]*"[^>]*/>', "", r)
        with open(path, "w", encoding="utf8") as fh:
            fh.write(r)
    with open(ct, encoding="utf8") as fh:
        types = fh.read()
    types = re.sub(r'<Override PartName="/ppt/(slides|notesSlides)/[^"]*"[^>]*/>', "", types)
    with open(ct, "w", encoding="utf8") as fh:
        fh.write(types)
    out = os.path.join(work, "_base.pptx")
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(work):
            for name in files:
                full = os.path.join(root, name)
                if full == out:
                    continue
                z.write(full, os.path.relpath(full, work))
    return out


TEMPLATE = sys.argv[1] if len(sys.argv) > 1 else "Acc_PPT_IMP_StarterPack.potx"
OUTPUT = sys.argv[2] if len(sys.argv) > 2 else "acn_jira_registration_timing.pptx"

# ---------------------------------------------------------------- palette ----
PURPLE      = RGBColor(0xA1, 0x00, 0xFF)   # accent1 - Accenture Purple
PURPLE_MID  = RGBColor(0x75, 0x00, 0xC0)   # accent2
PURPLE_DARK = RGBColor(0x46, 0x00, 0x73)   # accent3
PURPLE_LT   = RGBColor(0xBE, 0x82, 0xFF)   # accent5
PURPLE_PALE = RGBColor(0xDC, 0xAF, 0xFF)   # accent6
TINT        = RGBColor(0xF4, 0xE9, 0xFF)   # pale purple wash
WARM        = RGBColor(0xE6, 0xE6, 0xDC)   # lt2
WARM_GREY   = RGBColor(0x96, 0x96, 0x8C)   # dk2
BLACK       = RGBColor(0x00, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x33, 0x33, 0x33)
MUTED       = RGBColor(0x70, 0x70, 0x6A)

JP = "Meiryo"      # East-Asian face
LAT = "Arial"      # Latin face

# ------------------------------------------------------------------ grid -----
L, R = 0.42, 12.92          # ACN content margins
TRACK_L = 1.62              # swim-lane track starts after the row labels
N = 9
PITCH = (R - TRACK_L) / N
CHEV_W = PITCH - 0.045
def cx(i):  return TRACK_L + PITCH * i + PITCH / 2      # centre of phase i
def bx(i):  return TRACK_L + PITCH * i                  # left edge of phase i

SOP_X   = bx(6)          # boundary between 検証 and 量産
INF_X   = cx(2)          # 開発企画書 発行 -> INF trigger
INIT_X  = cx(3)          # 仕様FIX      -> イニシ trigger


# --------------------------------------------------------------- helpers ----
def style(run, size, bold=False, color=BLACK, face=LAT, ea=JP):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, face
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", ea if tag == "a:ea" else face)


def textbox(sl, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tb, tf


def line(sl, x1, y1, x2, y2, color, width=1.0, dash=None):
    ln = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        ln.line._get_or_add_ln().append(
            ln.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": dash}))
    return ln


def shape(sl, kind, x, y, w, h, fill=None, outline=None, ow=1.0):
    sh = sl.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if outline is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = outline; sh.line.width = Pt(ow)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return sh


def no_shadow(sh):
    sh.shadow.inherit = False


# ------------------------------------------------------------------ build ----
prs = Presentation(base_from_template(TEMPLATE))
layout = next(l for l in prs.slide_master.slide_layouts if l.name == "Long Headline Only")
sl = prs.slides.add_slide(layout)

# --- eyebrow -----------------------------------------------------------------
_, tf = textbox(sl, L, 0.36, 9.0, 0.24)
style(tf.paragraphs[0].add_run(), 11.5, True, PURPLE)
tf.paragraphs[0].runs[0].text = "JIRA登録タイミング ｜ INF・イニシ別の登録トリガー"

# --- headline (the layout's own title placeholder) ---------------------------
title = sl.shapes.title
title.left, title.top = Inches(L), Inches(0.62)
title.width, title.height = Inches(12.50), Inches(0.66)
ttf = title.text_frame
ttf.word_wrap = True
ttf.margin_left = ttf.margin_right = ttf.margin_top = ttf.margin_bottom = 0
ttf.vertical_anchor = MSO_ANCHOR.TOP
p = ttf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
style(p.add_run(), 27, True, BLACK, face="Arial Black")
p.runs[0].text = "INFは開発企画書発行時、イニシは仕様FIX時にJIRAへ登録する"

# --- deck (sub-headline) -----------------------------------------------------
_, tf = textbox(sl, L, 1.34, 12.0, 0.28)
style(tf.paragraphs[0].add_run(), 13.5, False, MUTED)
tf.paragraphs[0].runs[0].text = "登録判断の位置を入口・出口に分け、SOPを境界に管理区間を切り替える"

# --- definition band ---------------------------------------------------------
band = shape(sl, MSO_SHAPE.RECTANGLE, L, 1.78, R - L, 0.70, fill=TINT)
band.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
_, tf = textbox(sl, L + 0.22, 1.90, 11.9, 0.24)
p = tf.paragraphs[0]
for t, sz, b, c in [("登録判断", 12, True, PURPLE_DARK), ("     ", 12, False, BLACK),
                    ("INF ＝ 開発企画書 発行時", 12.5, True, BLACK),
                    ("     ｜     ", 12.5, False, WARM_GREY),
                    ("イニシ ＝ 仕様FIX時", 12.5, True, BLACK),
                    ("     ｜     ", 12.5, False, WARM_GREY),
                    ("登録以降はJIRAで進捗・課題を管理", 12.5, True, PURPLE_MID)]:
    r = p.add_run(); r.text = t; style(r, sz, b, c)
_, tf = textbox(sl, L + 0.22, 2.18, 11.9, 0.22)
style(tf.paragraphs[0].add_run(), 10.5, False, MUTED)
tf.paragraphs[0].runs[0].text = "登録前は個別に管理し、登録後は案件の進捗・課題をJIRA上で追跡する。"

# --- lifecycle label + SOP zone captions -------------------------------------
_, tf = textbox(sl, L, 2.66, 4.0, 0.24)
style(tf.paragraphs[0].add_run(), 11.5, True, PURPLE)
tf.paragraphs[0].runs[0].text = "プロジェクトライフサイクル"

for txt, x0, x1, col in [("SOP前", bx(4), bx(6), PURPLE_MID),
                         ("SOP",  SOP_X - 0.30, SOP_X + 0.30, PURPLE_DARK),
                         ("SOP後", bx(6), bx(8), PURPLE_DARK)]:
    _, tf = textbox(sl, x0, 2.64, x1 - x0, 0.24, align=PP_ALIGN.CENTER)
    style(tf.paragraphs[0].add_run(), 10.5, True, col)
    tf.paragraphs[0].runs[0].text = txt

# --- SOP boundary ------------------------------------------------------------
line(sl, SOP_X, 2.92, SOP_X, 5.44, PURPLE_DARK, 1.25, dash="dash")

# --- phase chevrons ----------------------------------------------------------
PHASES = [
    ("引合",     WARM,  MUTED,       None),
    ("提案",     WARM,  MUTED,       None),
    ("仕様検討", TINT,  PURPLE_DARK, PURPLE),
    ("投資",     TINT,  PURPLE_DARK, PURPLE),
    ("量産準備", TINT,  PURPLE_DARK, PURPLE),
    ("検証",     TINT,  PURPLE_DARK, PURPLE),
    ("量産",     PURPLE,      WHITE, None),
    ("補給品",   PURPLE_MID,  WHITE, None),
    ("EOP",      WARM,  MUTED,       None),
]
for i, (name, fill, fg, outline) in enumerate(PHASES):
    ch = shape(sl, MSO_SHAPE.CHEVRON, bx(i), 3.00, CHEV_W, 0.44,
               fill=fill, outline=outline, ow=1.0)
    style(ch.text_frame.paragraphs[0].add_run(), 11, True, fg)
    ch.text_frame.paragraphs[0].runs[0].text = name

# --- milestone axis + trigger markers ----------------------------------------
AXIS_Y = 3.92
for x, label in [(INF_X, "開発企画書 発行"), (INIT_X, "仕様FIX")]:
    _, tf = textbox(sl, x - 1.05, 3.56, 2.10, 0.22, align=PP_ALIGN.CENTER)
    style(tf.paragraphs[0].add_run(), 10.5, True, PURPLE_DARK)
    tf.paragraphs[0].runs[0].text = label
line(sl, TRACK_L, AXIS_Y, R, AXIS_Y, WARM_GREY, 0.75)
for x in (INF_X, INIT_X):
    shape(sl, MSO_SHAPE.DIAMOND, x - 0.075, AXIS_Y - 0.075, 0.15, 0.15, fill=PURPLE)

# --- swim lanes --------------------------------------------------------------
LANES = [("INF",   4.64, INF_X),
         ("イニシ", 5.26, INIT_X)]
for name, y, trig in LANES:
    # row label
    _, tf = textbox(sl, L, y - 0.13, 1.05, 0.26, align=PP_ALIGN.RIGHT,
                    anchor=MSO_ANCHOR.MIDDLE)
    style(tf.paragraphs[0].add_run(), 14, True, BLACK, face="Arial Black")
    tf.paragraphs[0].runs[0].text = name
    # track: grey before the trigger, purple after it
    line(sl, TRACK_L, y, trig, y, WARM_GREY, 1.25)
    line(sl, trig, y, R, y, PURPLE, 1.75)
    # trigger drop-line from the milestone axis into this lane
    line(sl, trig, AXIS_Y, trig, y, PURPLE_LT, 0.75, dash="sysDot")
    d = shape(sl, MSO_SHAPE.DIAMOND, trig - 0.075, y - 0.075, 0.15, 0.15, fill=PURPLE)
    # segment captions
    for txt, x0, x1, sz, b, col in [
            ("個別管理",     TRACK_L + 0.06, trig,  10.5, False, MUTED),
            ("JIRAで管理",   trig,          SOP_X,  11.5, True,  PURPLE_MID),
            ("PCR都度起票",  SOP_X,         R,      11.5, True,  PURPLE_DARK)]:
        _, tf = textbox(sl, x0, y - 0.34, x1 - x0, 0.22, align=PP_ALIGN.CENTER)
        style(tf.paragraphs[0].add_run(), sz, b, col)
        tf.paragraphs[0].runs[0].text = txt

# --- open items --------------------------------------------------------------
box = shape(sl, MSO_SHAPE.RECTANGLE, L, 5.62, R - L, 0.62, fill=WARM)
_, tf = textbox(sl, L + 0.22, 5.74, 1.4, 0.22)
style(tf.paragraphs[0].add_run(), 11.5, True, PURPLE_DARK)
tf.paragraphs[0].runs[0].text = "未決事項"
for txt, x0 in [("① 登録判断は全件か選別か（タイミング定義と判定基準の要否）", 1.90),
                ("② 補給品の管理媒体（JIRA／Excel）と既存方針の整合",      7.55)]:
    _, tf = textbox(sl, x0, 5.74, 5.35, 0.40)
    style(tf.paragraphs[0].add_run(), 11, False, INK)
    tf.paragraphs[0].runs[0].text = txt

# --- So What -----------------------------------------------------------------
sw = shape(sl, MSO_SHAPE.RECTANGLE, L, 6.30, R - L, 0.62, fill=PURPLE)
_, tf = textbox(sl, L + 0.30, 6.42, 1.6, 0.26, anchor=MSO_ANCHOR.MIDDLE)
style(tf.paragraphs[0].add_run(), 13, True, WHITE, face="Arial Black")
tf.paragraphs[0].runs[0].text = "So What"
_, tf = textbox(sl, L + 2.05, 6.42, 10.3, 0.26, anchor=MSO_ANCHOR.MIDDLE)
style(tf.paragraphs[0].add_run(), 13.5, True, WHITE)
tf.paragraphs[0].runs[0].text = "INF／イニシを登録起点として定義し、SOP前後の管理責任を明確にする"

# --- footer (ACN master geometry) --------------------------------------------
_, tf = textbox(sl, L + 0.32, 7.10, 6.0, 0.22)
style(tf.paragraphs[0].add_run(), 9, False, MUTED)
tf.paragraphs[0].runs[0].text = "JIRA登録タイミング ｜ トリガー定義"
_, tf = textbox(sl, 12.30, 7.10, 0.62, 0.22, align=PP_ALIGN.RIGHT)
style(tf.paragraphs[0].add_run(), 9, False, MUTED)
tf.paragraphs[0].runs[0].text = "8"

sl.notes_slide.notes_text_frame.text = (
    "INFは開発企画書の発行時点、イニシは仕様FIX時点をJIRA登録のトリガーとする。"
    "登録前は個別管理、登録後はJIRA上で進捗・課題を追跡し、SOPを境界に量産・補給品は"
    "PCRを都度起票する運用へ切り替える。未決は登録判断の全件／選別の別と、"
    "補給品の管理媒体（JIRA／Excel）の整合。")

prs.save(OUTPUT)
print("saved " + OUTPUT)
