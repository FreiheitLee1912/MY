# -*- coding: utf-8 -*-
"""ポテンシャル案件の生準管理ルール — ダイセル標準テンプレートで作成.

Usage:
    python build_deck_potential.py --template 2026_Standard_Template_WideScreenEN.pptx

版面・配色・書体は daicel_kit（テンプレートからの実測値、Meiryo）に従う。
"""
import argparse

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from titus_kit import BLACK, WHITE, style, say, rich, cell
from daicel_kit import (BLUE, DARK, PALE, GREY, RULE, MUTED, RED, GREEN,
                        TAG_DECIDED, TAG_OPEN, NOTE_FILL, WARN_FILL,
                        L, R, W, BODY_TOP, base_from_template, secbar, note, page)

TRIAL = ("試行運用", TAG_DECIDED)
TBD = ("要確定", TAG_OPEN)


def tcell(sl, x, y, w, h, text, size=11.5, bold=False, fg=BLACK, fill=WHITE,
          align=PP_ALIGN.LEFT, pad=0.14, space=None):
    return cell(sl, x, y, w, h, text, size, bold, fg, fill, RULE, 0.75,
                align=align, pad=pad, space=space)


def thead(sl, x, y, cols, h=0.36, size=11.5):
    for t, w in cols:
        cell(sl, x, y, w, h, t, size, True, WHITE, DARK, RULE, 0.75)
        x += w


def rows_table(sl, y, cols, rows, rh, key_fill=PALE, size=11, space=1.25):
    """Header + body rows. cols = [(label, width)]; rows = [(cell, ...)]."""
    thead(sl, L, y, cols)
    ry = y + 0.36
    for row in rows:
        x = L
        for i, ((_, w), val) in enumerate(zip(cols, row)):
            tcell(sl, x, ry, w, rh, val, size, i == 0,
                  DARK if i == 0 else BLACK, key_fill if i == 0 else WHITE,
                  PP_ALIGN.CENTER if i == 0 else PP_ALIGN.LEFT, space=space)
            x += w
        ry += rh
    return ry


# ------------------------------------------------------------------ 1 表紙 ----
def cover(prs):
    sl = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
    tf = sl.shapes.title.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "ポテンシャル案件の\n生準管理ルール"
    style(r, 26, True, WHITE)
    stf = sl.placeholders[1].text_frame
    stf.word_wrap = True
    for i, txt in enumerate(["生産準備部門　運用ルール",
                             "案（試行運用）　2026年9月4日"]):
        pp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        rr = pp.add_run(); rr.text = txt
        style(rr, 14 if i == 0 else 12, i == 0, WHITE)


# ---------------------------------------------------------------- 2 全体像 ----
def overview(prs, n):
    sl = page(prs, "本ルールの全体像",
              "Action Plan Rankと生準影響度をもとに、生準がどの案件を、どの粒度で、\n"
              "どのツールで見るかを整理したもの。", (TRIAL,), n)
    y = BODY_TOP
    cols = [("節", 0.80), ("表題", 3.10), ("要点", W - 3.90)]
    rows = [("1", "目的", "生準がモニタリングする対象はC+に限定する"),
            ("2", "管理の境界", "C+はkintone、A／Award以降はJIRAで正式管理する"),
            ("3", "3つの判定Gate", "管理対象・粒度・ツールをGateで切り替える"),
            ("4", "kintone登録項目", "C+案件サマリーに登録する項目を定める"),
            ("5", "生準が更新する項目", "生産準備に関する項目は生準が更新する"),
            ("6", "レビューと日常管理", "レビューの場を設けるのは影響度「大」のみ"),
            ("7", "共通ルール", "受注確定前にJIRA Taskを作成しない")]
    ry = rows_table(sl, y, cols, rows, 0.50, size=11.5)
    note(sl, L, ry + 0.26, W, 0.62, "判断の軸",
         "「どの案件を見るか」＝第2節　／　「どの粒度で見るか」＝第3・6節　／　"
         "「どのツールで見るか」＝第2・4節")


# -------------------------------------------------------------- 3 [1] 目的 ----
def purpose(prs, n):
    sl = page(prs, "[1] 目的",
              "ポテンシャル案件に対して生準がどこまで関与するかを整理し、"
              "案件の確度に応じて管理の手間に濃淡をつける。", (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "管理対象の限定")
    cell(sl, L, y + 0.38, W, 1.05,
         "全案件を一律に見るのではなく、生準がモニタリングする対象はC+に限定する。\n"
         "C+はMTP・Budget・投資判断に含まれるため、未受注であっても先行して状況を把握する。",
         14, True, DARK, NOTE_FILL, None, space=1.5)
    ny = y + 1.67
    secbar(sl, L, ny, W, "ねらい")
    goals = [("01", "受注前に工数をかけすぎない", "受注前の案件に過大な工数をかけない。"),
             ("02", "先行着手を要する案件を早期に掴む",
              "設備・金型・長納期品など、先行着手を要する案件を早期に把握する。"),
             ("03", "受注後は速やかにJIRAへ渡す",
              "受注確定後は速やかにJIRAへ移行し、正式な量産準備管理に接続する。")]
    ry = ny + 0.38
    for no, t, d in goals:
        cell(sl, L, ry, 0.70, 0.72, no, 12.5, True, WHITE, BLUE, RULE, 0.75)
        tcell(sl, L + 0.70, ry, 3.90, 0.72, t, 12, True, DARK, PALE)
        tcell(sl, L + 4.60, ry, W - 4.60, 0.72, d, 11.5)
        ry += 0.72


# ---------------------------------------------------------- 4 [2] 管理の境界 ----
def boundary(prs, n):
    sl = page(prs, "[2] 管理の境界",
              "プロダクトライフサイクル上、生準の管理はC+から始まり、"
              "A／Awardで正式管理へ切り替わる。", (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "プロダクトライフサイクルと管理の境界")
    phases = ["引き合い\nコンペ", "C+", "A／Award", "投資・設計",
              "量産準備", "PPAP\n顧客承認", "SOP", "PCR"]
    gap = 0.06
    pw = (W - gap * (len(phases) - 1)) / len(phases)
    bx = lambda i: L + i * (pw + gap)
    for i, name in enumerate(phases):
        hot = (i == 1)
        cell(sl, bx(i), y + 0.52, pw, 0.72, name, 11.5, hot,
             WHITE if hot else BLACK, BLUE if hot else GREY,
             RULE, 0.75, align=PP_ALIGN.CENTER, space=1.25)
    cell(sl, bx(1), y + 1.36, bx(3) - bx(1) - gap, 0.42,
         "生準モニタリング｜kintone", 11.5, True, WHITE, BLUE, None)
    cell(sl, bx(2), y + 1.86, bx(8) - bx(2) - gap, 0.42,
         "生準正式管理｜A／Award以降・JIRA新規立上げ", 11.5, True, WHITE, DARK, None)
    say(sl, L, y + 2.36, W, 0.24,
        "モニタリング対象はC+のみ（影響度別に管理）。"
        "Bランクは既存のAction Plan／MTP・Budget管理へ引き継ぐ。",
        11, False, MUTED, align=PP_ALIGN.RIGHT)
    ny = y + 2.72
    secbar(sl, L, ny, W, "C+の定義")
    cell(sl, L, ny + 0.38, W, 1.10,
         "未受注だが戦略的に獲得を前提として取り組んでおり、投資判断にも考慮されるターゲット案件。\n"
         "Opportunities not yet awarded but strategically pursued with intent to win "
         "and considered for investment decisions.",
         11.5, False, BLACK, NOTE_FILL, None, align=PP_ALIGN.CENTER, space=1.45)


# ------------------------------------------------------------ 5 [2] 表1 Rank ----
def ranks(prs, n):
    sl = page(prs, "[2] Potential Rankの定義と生準の関与",
              "生準がモニタリングするのはC+のみ。A・Bは既存管理、C以下は対象外。",
              (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表1　Potential Rankの定義と生準の関与")
    cols = [("Rank", 1.05), ("定義", 4.30), ("MTP・Budget・投資", 2.15), ("生準の関与", W - 7.50)]
    rows = [("A", "現行商権／受注獲得済のプログラム", "対象", "Award以降、JIRAで正式管理"),
            ("B", "現行商権の次期車／商権獲得が確定的なもの", "対象",
             "既存のAction Plan／MTP・Budget管理へ引き継ぐ"),
            ("C+", "未受注だが戦略的に獲得を前提として取り組んでおり、投資判断にも考慮されるターゲット案件",
             "対象", "kintoneでモニタリング（本ルールの対象）"),
            ("C", "営業としてのターゲットプログラムだが、Moduleとは合意前のもの", "対象外", "情報確認のみ"),
            ("D", "受注ターゲットにできるか、案件として先ずはリストに載せたもの", "対象外", "モニタリング対象外"),
            ("X", "現行商権の次期車でも商権を失ったもの", "対象外", "モニタリング対象外")]
    thead(sl, L, y + 0.38, cols, size=10.5)
    ry = y + 0.74
    for rank, dfn, mtp, role in rows:
        hot = rank == "C+"
        tcell(sl, L, ry, 1.05, 0.62, rank, 13, True,
              WHITE if hot else DARK, BLUE if hot else PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.05, ry, 4.30, 0.62, dfn, 10.5,
              fill=NOTE_FILL if hot else WHITE, space=1.25)
        tcell(sl, L + 5.35, ry, 2.15, 0.62, mtp, 11, mtp == "対象",
              DARK if mtp == "対象" else MUTED,
              NOTE_FILL if hot else WHITE, PP_ALIGN.CENTER)
        tcell(sl, L + 7.50, ry, W - 7.50, 0.62, role, 10.5, hot,
              BLACK, NOTE_FILL if hot else WHITE, space=1.25)
        ry += 0.62


# ------------------------------------------------------------- 6 [3] Gate ----
def gates(prs, n):
    sl = page(prs, "[3] 3つの判定Gate",
              "管理対象・粒度・ツールは、次の3つのGateで切り替える。", (TRIAL, TBD), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "判定Gate")
    g = [("01", "C+該当／管理開始", BLUE,
          "Action PlanでC+となった案件を、生準が継続して見るかどうかを判断する。",
          "Action PlanでC+と判定／想定SOP・顧客判断時期／全体管理対象（要確定）",
          "判定後　C+＝kintone登録"),
         ("02", "C+生準影響度判定", DARK,
          "生準影響を小・中・大に分け、管理の濃淡と対応の優先順位を決める。",
          "設備・金型／工程／能力／拠点／長納期品／日程・技術・品質／投資／"
          "想定SOP時期／新規品種か否か",
          "判定後　小＝簡易／中＝標準／大＝重点"),
         ("03", "正式案件化", BLUE,
          "受注後の量産準備をJIRAで正式に管理する。",
          "Nomination／採用内示／正式発注など、受注・採用の確定",
          "判定後　JIRAへ新規立上げ案件を登録")]
    cw = (W - 0.40) / 3
    gy = y + 0.50
    for i, (no, t, col, aim, cond, res) in enumerate(g):
        x = L + i * (cw + 0.20)
        cell(sl, x, gy, cw, 0.44, f"{no}　{t}", 12.5, True, WHITE, col, None)
        tcell(sl, x, gy + 0.44, cw, 0.22, "目的", 9.5, True, MUTED, PALE, pad=0.12)
        tcell(sl, x, gy + 0.66, cw, 0.62, aim, 10.5, pad=0.12, space=1.3)
        tcell(sl, x, gy + 1.28, cw, 0.22, "主な判定条件", 9.5, True, MUTED, PALE, pad=0.12)
        tcell(sl, x, gy + 1.50, cw, 0.86, cond, 10.5, pad=0.12, space=1.3)
        cell(sl, x, gy + 2.36, cw, 0.40, res, 11, True, col, PALE, RULE, 0.75)
    ny = gy + 2.98
    nw = (W - 0.20) / 2
    note(sl, L, ny, nw, 0.88, "優先順位の考え方",
         "想定SOPが早い案件と新規品種の案件を優先して対応する。", size=11)
    note(sl, L + nw + 0.20, ny, nw, 0.88, "作業依頼が発生した場合",
         "ATP上のポテンシャルをC+に変更したうえで依頼する。Rankを変えずに依頼しない。",
         fill=WARN_FILL, size=11)


# ------------------------------------------------------- 7 [4] kintone項目 ----
def kintone(prs, n):
    sl = page(prs, "[4] kintone登録項目",
              "C+と判定された案件は、kintoneのC+案件サマリーに次の項目を登録する。",
              (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表2　C+案件サマリーの登録項目")
    cols = [("区分", 2.30), ("項目", W - 2.30)]
    rows = [("案件基本", "案件ID／顧客／Project・車種／対象製品"),
            ("Rank", "C+／C+判定日"),
            ("担当", "生準担当"),
            ("日程・規模", "想定SOP／顧客判断日／数量／拠点／次回確認日"),
            ("生準影響", "小・中・大／設備・金型・工程・能力"),
            ("Action・結果", "次アクション／先行着手承認／案件結果")]
    ry = rows_table(sl, y + 0.38, cols, rows, 0.56, size=12)
    note(sl, L, ry + 0.28, W, 0.72, "登録の起点",
         "Gate 01でC+と判定した時点で登録する。C+以外の案件はkintoneに登録しない（第7節）。")


# ------------------------------------------------------ 8 [5] 生準更新項目 ----
def updates(prs, n):
    sl = page(prs, "[5] 生準が更新する項目",
              "kintone上の生産準備に関する項目は、生準が更新する。生準以外の項目は書き換えない。",
              (TBD,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表3　生産準備項目（ドラフト・要精査）")
    cols = [("区分", 1.75), ("項目", 6.30), ("更新の目安", W - 8.05)]
    rows = [("判定・区分", "生準影響度（小・中・大）／影響度判定日／優先順位", "Gate 02判定時、区分変更時"),
            ("生産条件", "想定生産拠点／新規・流用の別／対象製品", "初回登録時、変更時"),
            ("設備・金型", "設備（新設・改造・流用）／金型（新規・改造・流用、面数）／必要投資の概算",
             "見込みが立った時点"),
            ("工程・能力", "新規工程の要否／必要能力／現有能力との差／能力確保の見通し",
             "見込みが立った時点"),
            ("日程", "想定SOPからの逆算日程（設備手配・型手配・トライ・量産試作）／長納期品の有無と品目",
             "想定SOP変更時"),
            ("先行対応", "先行着手の要否／承認状況／着手日", "判断・承認の都度"),
            ("課題・Action", "技術・品質課題／リスクと対応状況／次アクションと期限／次回確認日",
             "状況が動いた都度")]
    ry = rows_table(sl, y + 0.38, cols, rows, 0.48, size=10.5, space=1.2)
    note(sl, L, ry + 0.20, W, 0.62, "この表はドラフト",
         "kintoneの実際の項目に合わせて過不足を精査する。特に投資概算と逆算日程は、"
         "どこまでC+段階で記入できるか運用しながら決める。", fill=WARN_FILL, size=11)


# ---------------------------------------------------------- 9 [6] レビュー ----
def review(prs, n):
    sl = page(prs, "[6] レビューと日常管理",
              "レビューの場を設けるのは影響度「大」のみ。小・中は個別担当が管理し、"
              "変化点が出たときに共有する。", (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表4　影響度別の管理方法")
    cols = [("影響度", 1.60), ("管理方法", 3.60), ("頻度", 2.60), ("見る内容", W - 7.80)]
    thead(sl, L, y + 0.38, cols)
    rows = [("小", "個別担当で管理（レビューなし）", "随時", "変化点のみ", False),
            ("中", "個別担当で管理（レビューなし）", "随時", "影響・Risk・Action", False),
            ("大", "レビューを実施", "隔週または重要マイルストーン",
             "投資・長納期品・逆算日程", True)]
    ry = y + 0.74
    for lvl, how, freq, what, hot in rows:
        tcell(sl, L, ry, 1.60, 0.70, lvl, 15, True,
              WHITE if hot else DARK, BLUE if hot else PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.60, ry, 3.60, 0.70, how, 11.5, hot, BLACK,
              NOTE_FILL if hot else WHITE)
        tcell(sl, L + 5.20, ry, 2.60, 0.70, freq, 11.5, hot, BLACK,
              NOTE_FILL if hot else WHITE)
        tcell(sl, L + 7.80, ry, W - 7.80, 0.70, what, 11.5, hot, BLACK,
              NOTE_FILL if hot else WHITE)
        ry += 0.70
    note(sl, L, ry + 0.30, W, 0.86, "運用",
         "小・中はレビューの場を設けない。変化点が出たときに担当から共有する。"
         "「大」は隔週または重要マイルストーンで、投資・長納期品・逆算日程を見る。")


# -------------------------------------------------------- 10 [7] 共通ルール ----
def common(prs, n):
    sl = page(prs, "[7] 共通ルール",
              "確度が上がる前に管理を始めない。受注確定をもってJIRAへ引き渡す。",
              (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "やらないこと", fill=RED)
    dont = [("RFQ・提案・コンペだけを理由に生準管理を始めない",
             "情報の確認にとどめる。"),
            ("C+以外の案件をkintoneに登録しない",
             "登録対象はGate 01でC+と判定した案件のみ。"),
            ("受注確定前にJIRA Taskを作成しない",
             "JIRAへの登録はGate 03（Nomination／採用内示／正式発注）以降。")]
    ry = y + 0.38
    for i, (t, d) in enumerate(dont, 1):
        cell(sl, L, ry, 0.70, 0.66, "✕", 15, True, WHITE, RED, RULE, 0.75)
        tcell(sl, L + 0.70, ry, 5.30, 0.66, t, 12, True, RED, PALE)
        tcell(sl, L + 6.00, ry, W - 6.00, 0.66, d, 11.5)
        ry += 0.66
    ny = ry + 0.32
    secbar(sl, L, ny, W, "受注確定後")
    cell(sl, L, ny + 0.38, W, 0.86,
         "受注確定後、JIRAでの記録の残し方は「JIRA文書管理規程」による。",
         14, True, DARK, NOTE_FILL, None)


# ---------------------------------------------------------------- 11 まとめ ----
def closing(prs, n):
    sl = page(prs, "まとめ",
              "本ルールの要点と、試行にあたって残っている確定事項。", (TRIAL, TBD), n)
    y = BODY_TOP
    cell(sl, L, y, W, 0.78,
         "見るのはC+だけ。粒度は影響度で決め、受注確定でJIRAへ渡す",
         20, True, WHITE, BLUE, None)
    pts = [("どの案件を見るか", "モニタリングはC+のみ。\nA・Bは既存管理、C以下は対象外。"),
           ("どの粒度で見るか", "影響度で小・中・大に分け、\nレビューは「大」のみ。"),
           ("どのツールで見るか", "C+はkintone、A／Award以降は\nJIRAで正式管理。")]
    cw = (W - 0.40) / 3
    for i, (t, d) in enumerate(pts):
        x = L + i * (cw + 0.20)
        cell(sl, x, y + 1.04, cw, 0.40, t, 12.5, True, WHITE, DARK, None)
        tcell(sl, x, y + 1.44, cw, 1.00, d, 11.5, align=PP_ALIGN.CENTER, space=1.35)
    ny = y + 2.76
    secbar(sl, L, ny, W, "確定が必要な事項")
    opens = [("01", "Gate 01の判定条件「全体管理対象」の定義"),
             ("02", "表3（生産準備項目）の精査 ― 投資概算・逆算日程をC+段階でどこまで記入するか"),
             ("03", "本ルールは案。試行しながら実態に合わせて見直す")]
    ry = ny + 0.38
    for no, t in opens:
        tcell(sl, L, ry, 0.66, 0.44, no, 11, True, BLACK, TAG_OPEN, PP_ALIGN.CENTER)
        tcell(sl, L + 0.66, ry, W - 0.66, 0.44, t, 11.5, pad=0.16)
        ry += 0.44


def back_cover(prs):
    sl = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
    say(sl, 3.82, 2.73, 5.68, 2.61, "Thank you", 32, True, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return sl




# ============================================================================
# 3ページ版 — 12枚を「どの案件を／どの粒度で／何を記録するか」の3軸に畳んだもの
# ============================================================================

PHASES = ["引き合い\nコンペ", "C+", "A／Award", "投資・設計",
          "量産準備", "PPAP\n顧客承認", "SOP", "PCR"]


def p1_scope(prs, n):
    """どの案件を見るか — 管理の境界とRank定義."""
    sl = page(prs, "[1-2] 管理対象と境界 ― 見るのはC+だけ",
              "生準がモニタリングする対象はC+に限定する。C+はMTP・Budget・投資判断に"
              "含まれるため、未受注でも先行して状況を把握する。", (TRIAL,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "プロダクトライフサイクルと管理の境界", h=0.32, size=12.5)
    gap = 0.06
    pw = (W - gap * (len(PHASES) - 1)) / len(PHASES)
    bx = lambda i: L + i * (pw + gap)
    for i, name in enumerate(PHASES):
        hot = (i == 1)
        cell(sl, bx(i), y + 0.40, pw, 0.56, name, 10.5, hot,
             WHITE if hot else BLACK, BLUE if hot else GREY,
             RULE, 0.75, align=PP_ALIGN.CENTER, space=1.2)
    cell(sl, bx(1), y + 1.02, bx(3) - bx(1) - gap, 0.34,
         "生準モニタリング｜kintone", 10.5, True, WHITE, BLUE, None)
    cell(sl, bx(2), y + 1.40, bx(8) - bx(2) - gap, 0.34,
         "生準正式管理｜A／Award以降・JIRA新規立上げ", 10.5, True, WHITE, DARK, None)

    ty = y + 1.96
    secbar(sl, L, ty, W, "表1　Potential Rankの定義と生準の関与", h=0.32, size=12.5)
    cols = [("Rank", 0.90), ("定義", 4.55), ("MTP・Budget・投資", 1.75), ("生準の関与", W - 7.20)]
    x = L
    for t, w in cols:
        cell(sl, x, ty + 0.32, w, 0.30, t, 10, True, WHITE, DARK, RULE, 0.75)
        x += w
    rows = [("A", "現行商権／受注獲得済のプログラム", "対象", "Award以降、JIRAで正式管理"),
            ("B", "現行商権の次期車／商権獲得が確定的なもの", "対象",
             "既存のAction Plan／MTP・Budget管理へ引き継ぐ"),
            ("C+", "未受注だが戦略的に獲得を前提として取り組んでおり、投資判断にも考慮されるターゲット案件",
             "対象", "kintoneでモニタリング（本ルールの対象）"),
            ("C", "営業としてのターゲットプログラムだが、Moduleとは合意前のもの", "対象外", "情報確認のみ"),
            ("D", "受注ターゲットにできるか、案件として先ずはリストに載せたもの", "対象外", "モニタリング対象外"),
            ("X", "現行商権の次期車でも商権を失ったもの", "対象外", "モニタリング対象外")]
    ry = ty + 0.62
    for rank, dfn, mtp, role in rows:
        hot = rank == "C+"
        tcell(sl, L, ry, 0.90, 0.40, rank, 11.5, True,
              WHITE if hot else DARK, BLUE if hot else PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 0.90, ry, 4.55, 0.40, dfn, 9.5,
              fill=NOTE_FILL if hot else WHITE, pad=0.10, space=1.15)
        tcell(sl, L + 5.45, ry, 1.75, 0.40, mtp, 10, mtp == "対象",
              DARK if mtp == "対象" else MUTED,
              NOTE_FILL if hot else WHITE, PP_ALIGN.CENTER)
        tcell(sl, L + 7.20, ry, W - 7.20, 0.40, role, 9.5, hot, BLACK,
              NOTE_FILL if hot else WHITE, pad=0.10, space=1.15)
        ry += 0.40


def p2_grain(prs, n):
    """どの粒度で見るか — 判定Gateと影響度別の管理."""
    sl = page(prs, "[3][6] 判定Gateと管理の濃淡",
              "管理対象・粒度・ツールは3つのGateで切り替え、影響度で管理の濃淡をつける。",
              (TRIAL, TBD), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "3つの判定Gate", h=0.32, size=12.5)
    g = [("01　C+該当／管理開始", BLUE,
          "Action PlanでC+と判定／想定SOP・顧客判断時期／全体管理対象（要確定）",
          "判定後　C+＝kintone登録"),
         ("02　C+生準影響度判定", DARK,
          "設備・金型／工程／能力／拠点／長納期品／日程・技術・品質／投資／"
          "想定SOP時期／新規品種か否か",
          "判定後　小＝簡易／中＝標準／大＝重点"),
         ("03　正式案件化", BLUE,
          "Nomination／採用内示／正式発注など、受注・採用の確定",
          "判定後　JIRAへ新規立上げ案件を登録")]
    cw = (W - 0.40) / 3
    gy = y + 0.40
    for i, (t, col, cond, res) in enumerate(g):
        x = L + i * (cw + 0.20)
        cell(sl, x, gy, cw, 0.38, t, 12, True, WHITE, col, None)
        tcell(sl, x, gy + 0.38, cw, 0.20, "主な判定条件", 9, True, MUTED, PALE, pad=0.10)
        tcell(sl, x, gy + 0.58, cw, 0.90, cond, 10, pad=0.12, space=1.28)
        cell(sl, x, gy + 1.48, cw, 0.36, res, 10.5, True, col, PALE, RULE, 0.75)

    ty = gy + 2.06
    secbar(sl, L, ty, W, "表4　影響度別の管理方法", h=0.32, size=12.5)
    cols = [("影響度", 1.50), ("管理方法", 3.50), ("頻度", 2.60), ("見る内容", W - 7.60)]
    x = L
    for t, w in cols:
        cell(sl, x, ty + 0.32, w, 0.30, t, 10, True, WHITE, DARK, RULE, 0.75)
        x += w
    rows = [("小", "個別担当で管理（レビューなし）", "随時", "変化点のみ", False),
            ("中", "個別担当で管理（レビューなし）", "随時", "影響・Risk・Action", False),
            ("大", "レビューを実施", "隔週または重要マイルストーン",
             "投資・長納期品・逆算日程", True)]
    ry = ty + 0.62
    for lvl, how, freq, what, hot in rows:
        tcell(sl, L, ry, 1.50, 0.46, lvl, 13, True,
              WHITE if hot else DARK, BLUE if hot else PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.50, ry, 3.50, 0.46, how, 10.5, hot, BLACK,
              NOTE_FILL if hot else WHITE, pad=0.12)
        tcell(sl, L + 5.00, ry, 2.60, 0.46, freq, 10.5, hot, BLACK,
              NOTE_FILL if hot else WHITE, pad=0.12)
        tcell(sl, L + 7.60, ry, W - 7.60, 0.46, what, 10.5, hot, BLACK,
              NOTE_FILL if hot else WHITE, pad=0.12)
        ry += 0.46
    note(sl, L, ry + 0.20, W, 0.44, "優先順位",
         "想定SOPが早い案件と新規品種の案件を優先する。"
         "生準に作業を依頼するときはATP上のポテンシャルをC+に変更したうえで依頼する。",
         size=10.5)


def p3_record(prs, n):
    """何を記録し、何をしないか — kintone項目と共通ルール."""
    sl = page(prs, "[4][5][7] kintone登録項目と共通ルール",
              "C+と判定した案件はkintoneに登録し、生産準備に関する項目は生準が更新する。",
              (TRIAL, TBD), n)
    y = BODY_TOP
    cw = (W - 0.20) / 2
    secbar(sl, L, y, cw, "表2　C+案件サマリーの登録項目", h=0.32, size=12)
    reg = [("案件基本", "案件ID／顧客／Project・車種／対象製品"),
           ("Rank", "C+／C+判定日"),
           ("担当", "生準担当"),
           ("日程・規模", "想定SOP／顧客判断日／数量／拠点／次回確認日"),
           ("生準影響", "小・中・大／設備・金型・工程・能力"),
           ("Action・結果", "次アクション／先行着手承認／案件結果")]
    ry = y + 0.32
    for k, v in reg:
        tcell(sl, L, ry, 1.55, 0.44, k, 10, True, DARK, PALE, PP_ALIGN.CENTER, pad=0.08)
        tcell(sl, L + 1.55, ry, cw - 1.55, 0.44, v, 9.5, pad=0.10, space=1.15)
        ry += 0.44

    x2 = L + cw + 0.20
    secbar(sl, x2, y, cw, "表3　生準が更新する項目（ドラフト）", h=0.32, size=12)
    upd = [("判定・区分", "生準影響度／影響度判定日／優先順位"),
           ("生産条件", "想定生産拠点／新規・流用の別／対象製品"),
           ("設備・金型", "設備・金型の新設／改造／流用／必要投資の概算"),
           ("工程・能力", "新規工程の要否／必要能力／現有能力との差"),
           ("日程", "想定SOPからの逆算日程／長納期品の有無と品目"),
           ("先行対応", "先行着手の要否／承認状況／着手日"),
           ("課題・Action", "技術・品質課題／リスク／次アクションと期限")]
    ry2 = y + 0.32
    for k, v in upd:
        tcell(sl, x2, ry2, 1.55, 0.377, k, 10, True, DARK, PALE, PP_ALIGN.CENTER, pad=0.08)
        tcell(sl, x2 + 1.55, ry2, cw - 1.55, 0.377, v, 9.5, pad=0.10, space=1.15)
        ry2 += 0.377

    ny = max(ry, ry2) + 0.22
    secbar(sl, L, ny, W, "やらないこと", h=0.32, size=12.5, fill=RED)
    dont = ["RFQ・提案・コンペだけを理由に\n生準管理を始めない",
            "C+以外の案件を\nkintoneに登録しない",
            "受注確定前に\nJIRA Taskを作成しない"]
    dw = (W - 0.40) / 3
    for i, t in enumerate(dont):
        cell(sl, L + i * (dw + 0.20), ny + 0.40, dw, 0.62, t, 11, True, RED,
             WARN_FILL, RULE, 0.75, align=PP_ALIGN.CENTER, space=1.3)
    note(sl, L, ny + 1.14, W, 0.46, "受注確定後・要確定",
         "受注確定後の記録は「JIRA文書管理規程」による。"
         "Gate 01の「全体管理対象」の定義と表3の項目精査は要確定。本ルールは案。",
         size=10.5)


def build_three(template, output):
    prs = Presentation(base_from_template(template))
    for i, fn in enumerate([p1_scope, p2_grain, p3_record], start=1):
        fn(prs, i)
    prs.save(output)
    print("saved", output, "-", len(prs.slides._sldIdLst), "slides")


def build(template, output):
    prs = Presentation(base_from_template(template))
    cover(prs)
    for i, fn in enumerate([overview, purpose, boundary, ranks, gates,
                            kintone, updates, review, common, closing], start=2):
        fn(prs, i)
    back_cover(prs)
    prs.save(output)
    print("saved", output, "-", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--pages", choices=["12", "3"], default="12",
                    help="12=詳細版 / 3=3ページ凝縮版")
    ap.add_argument("-o", "--output")
    a = ap.parse_args()
    if a.pages == "3":
        build_three(a.template, a.output or "potential_management_rule_3p.pptx")
    else:
        build(a.template, a.output or "potential_management_rule_daicel.pptx")
