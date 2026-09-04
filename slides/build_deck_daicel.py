# -*- coding: utf-8 -*-
"""JIRA文書管理規程 — ダイセル標準テンプレートで作成.

Usage:
    python build_deck_daicel.py --template 2026_Standard_Template_WideScreenEN.pptx

テンプレートのマスター（ロゴ・上部罫・下部著作権帯）はそのまま継承し、
本文だけを組む。配色はテンプレートのロゴから採ったブランドブルー。
"""
import argparse

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from titus_kit import BLACK, WHITE, style, say, rich, cell
from daicel_kit import (BLUE, DARK, PALE, GREY, RULE, MUTED, RED, GREEN,
                        TAG_DECIDED, TAG_OPEN, NOTE_FILL, WARN_FILL,
                        L, R, W, BODY_TOP, base_from_template, drop_placeholder,
                        secbar, note, foot, page)

DECIDED = ("決定事項", TAG_DECIDED)
OPEN = ("継続議論", TAG_OPEN)


def tcell(sl, x, y, w, h, text, size=12, bold=False, fg=BLACK, fill=WHITE,
          align=PP_ALIGN.LEFT, pad=0.14, space=None):
    return cell(sl, x, y, w, h, text, size, bold, fg, fill, RULE, 0.75,
                align=align, pad=pad, space=space)


def thead(sl, x, y, cols, h=0.36, size=12):
    for t, w in cols:
        cell(sl, x, y, w, h, t, size, True, WHITE, DARK, RULE, 0.75)
        x += w


# ------------------------------------------------------------------ 1 表紙 ----
def cover(prs):
    sl = prs.slides.add_slide(prs.slide_master.slide_layouts[0])
    t = sl.shapes.title
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "JIRA文書管理規程"
    style(r, 30, True, WHITE)
    sub = sl.placeholders[1]
    stf = sub.text_frame
    stf.word_wrap = True
    for i, (txt, sz, bold) in enumerate([("生産準備部門　運用規程", 15, True),
                                         ("第1.0版　2026年9月4日", 13, False)]):
        pp = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        rr = pp.add_run(); rr.text = txt
        style(rr, sz, bold, WHITE)


# ---------------------------------------------------------------- 2 全体像 ----
def overview(prs, n):
    sl = page(prs, "本規程の全体像と条文の構成",
              "本規程は、記録の保管先と登録可否を定める第1条〜第7条で構成する。\n"
              "第3条・第6条が登録可否の判断軸、第5条が保管先の切り分けにあたる。",
              (DECIDED,), n)
    y = BODY_TOP
    cols = [("条", 1.05), ("表題", 3.15), ("要点", W - 4.20)]
    thead(sl, L, y, cols)
    rows = [("第1条", "目的", "記録は個人のメールボックスではなく、当該タスクのチケットに保管する"),
            ("第2条", "適用範囲および公開範囲", "生準メンバのみ → 関連部署へ段階的に公開する"),
            ("第3条", "保管対象", "社外とのやり取りを最優先で登録し、図面・個人情報は登録しない"),
            ("第4条", "社外コミュニケーションの記録", "メールは原本のまま保管し、口頭決定は確認メールを送って残す"),
            ("第5条", "JIRAおよびWinchillの役割分担", "図面・CADと版管理はWinchillに寄せ、JIRAには参照リンクのみ"),
            ("第6条", "情報区分と取扱い", "機密以上はJIRAに登録せず、参照リンクのみを記載する"),
            ("第7条", "登録前チェックリスト", "登録前に5点を確認してから登録する")]
    ry = y + 0.36
    for no, title, pt in rows:
        tcell(sl, L, ry, 1.05, 0.50, no, 12, True, DARK, PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.05, ry, 3.15, 0.50, title, 12, True)
        tcell(sl, L + 4.20, ry, W - 4.20, 0.50, pt, 12)
        ry += 0.50
    note(sl, L, ry + 0.28, W, 0.62, "判断の軸",
         "「何を登録するか」＝第3条・第6条　／　「どこに保管するか」＝第5条　／　"
         "「どう残すか」＝第4条")


# ------------------------------------------------------------- 3 第1条 目的 ----
def purpose(prs, n):
    sl = page(prs, "[第1条] 目的",
              "JIRAは進捗管理の手段にとどまらず、「誰が・いつ・何を決定したか」を後から\n"
              "確認できる記録の保管先として運用する。", (DECIDED,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "記録を残す目的")
    goals = [("01", "認識の相違を防止する", "社外との協議・合意の経緯を\n確実に保全する。"),
             ("02", "追跡できる状態を保つ", "担当者の交代または不在時においても、\n第三者が経緯を追跡できる。"),
             ("03", "事実関係を再構成する", "問題発生時に、事実関係を速やかに\n再構成できるようにする。")]
    cw = (W - 0.40) / 3
    for i, (no, t, d) in enumerate(goals):
        x = L + i * (cw + 0.20)
        cell(sl, x, y + 0.58, cw, 0.40, f"{no}　{t}", 13, True, WHITE, DARK, None)
        tcell(sl, x, y + 0.98, cw, 1.42, d, 12, align=PP_ALIGN.CENTER, space=1.35)
    ny = y + 2.72
    secbar(sl, L, ny, W, "原則")
    cell(sl, L, ny + 0.38, W, 1.30,
         "記録は個人のメールボックスではなく、当該タスクのチケットに保管する。\n"
         "個人の受信箱のみに存在する情報は、組織の記録として扱わない。",
         15, True, DARK, NOTE_FILL, None, space=1.5)


# --------------------------------------------------------- 4 第2条 公開範囲 ----
def scope(prs, n):
    sl = page(prs, "[第2条] 適用範囲および公開範囲",
              "関連部署への公開を前提とするため、登録内容は当該範囲で閲覧されることを\n"
              "想定して記載する。", (DECIDED, OPEN), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表1　JIRAの公開範囲")
    cols = [("区分", 1.60), ("対象", 4.20), ("備考", W - 5.80)]
    thead(sl, L, y + 0.38, cols)
    rows = [("現行", "生準メンバのみ", "試行運用期間"),
            ("今後", "関連部署に公開する", "公開時期は別途通知する　※時期未定")]
    ry = y + 0.74
    for k, tgt, rem in rows:
        tcell(sl, L, ry, 1.60, 0.50, k, 12, True, DARK, PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.60, ry, 4.20, 0.50, tgt, 12)
        tcell(sl, L + 5.80, ry, W - 5.80, 0.50, rem, 12)
        ry += 0.50
    ny = ry + 0.30
    secbar(sl, L, ny, W, "公開拡大に備えて徹底すること")
    prep = [("社外情報は原本で残す", "要約ではなく原本を保管し、\n閲覧者が経緯を追える状態にする。"),
            ("個人情報を含めない", "業務上不要な連絡先・私的な内容は\n登録前に取り除く。"),
            ("機密はリンクのみ", "機密以上はWinchill等に保管し、\nJIRAには参照先だけを記載する。")]
    cw = (W - 0.40) / 3
    for i, (t, d) in enumerate(prep):
        x = L + i * (cw + 0.20)
        cell(sl, x, ny + 0.48, cw, 0.38, t, 12.5, True, WHITE, BLUE, None)
        tcell(sl, x, ny + 0.86, cw, 0.92, d, 11.5, align=PP_ALIGN.CENTER, space=1.32)
    note(sl, L, ny + 2.06, W, 0.58, "留意点",
         "関連部署に閲覧されて差し支えない表現か、登録の都度確認する（第7条のチェック項目に対応）。")


# --------------------------------------------------------- 5 第3条 保管対象 ----
def targets(prs, n):
    sl = page(prs, "[第3条] 保管対象",
              "プロジェクト管理に関連する以下の情報を、該当するチケットに保管する。",
              (DECIDED,), n)
    y = BODY_TOP
    cw = 5.65
    secbar(sl, L, y, cw, "登録するもの")
    rows = [("社外とのやり取り　※最重要", "顧客・サプライヤとの依頼、回答、条件調整、合意", True),
            ("決定事項", "方針決定および重要な決議内容", False),
            ("進捗報告", "定期報告、マイルストーンの達成状況", False),
            ("要件確認", "仕様決定、変更依頼の承認", False),
            ("議事録", "打合せの記録および議論の経過", False),
            ("ステータス", "課題・リスク・対応状況", False)]
    ry = y + 0.38
    for k, v, top in rows:
        tcell(sl, L, ry, 2.45, 0.62, k, 11.5, True, RED if top else DARK, PALE,
              pad=0.12, space=1.2)
        tcell(sl, L + 2.45, ry, cw - 2.45, 0.62, v, 11, pad=0.12, space=1.2)
        ry += 0.62
    x2 = L + cw + 0.20
    cw2 = W - cw - 0.20
    secbar(sl, x2, y, cw2, "登録しないもの", fill=RED)
    nos = [("図面・設計図・CADファイル", "Winchillにて保管する（第5条）"),
           ("業務上不要な個人情報", "個人携帯番号、自宅住所、私的な内容")]
    ry2 = y + 0.38
    for k, v in nos:
        cell(sl, x2, ry2, cw2, 0.40, k, 12.5, True, RED, PALE, RULE, 0.75,
             align=PP_ALIGN.LEFT, pad=0.16)
        tcell(sl, x2, ry2 + 0.40, cw2, 0.50, v, 11.5, pad=0.16)
        ry2 += 1.02
    note(sl, x2, ry2 + 0.14, cw2, 1.56, "判断に迷う場合",
         "登録前に第7条のチェックリストで確認する。機密区分に該当するものは、"
         "JIRAには参照リンクのみを記載する。", fill=WARN_FILL, size=11.5)


# ------------------------------------------ 6 第4条 社外コミュニケーション ----
def external(prs, n):
    sl = page(prs, "[第4条] 社外コミュニケーションの記録",
              "社外とのメールは原本のまま保管する。要約のみの記載は認めない。",
              (DECIDED,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "記録の原則")
    rules = [("01", "原本のまま保管する",
              "社外とのメールは、該当チケットに原本のまま保管する。要約のみの記載は認めない。"),
             ("02", "依頼と回答を対で保管する",
              "依頼と回答は対で保管し、合意内容を読み取れる状態とする。"),
             ("03", "訂正は追記で残す",
              "訂正が生じた場合、元の記録は削除せず、コメントにより追記する。")]
    ry = y + 0.38
    for no, t, d in rules:
        cell(sl, L, ry, 0.72, 0.72, no, 13, True, WHITE, BLUE, RULE, 0.75)
        tcell(sl, L + 0.72, ry, 3.48, 0.72, t, 12.5, True, DARK, PALE)
        tcell(sl, L + 4.20, ry, W - 4.20, 0.72, d, 11.5)
        ry += 0.72
    ny = ry + 0.30
    secbar(sl, L, ny, W, "電話・口頭で決定した場合", fill=DARK)
    cell(sl, L, ny + 0.38, W, 1.16,
         "内容の確認メールを相手方へ送付し、当該メールをチケットに保管する。\n"
         "本項の徹底が、認識の相違に対する最も有効な予防措置となる。",
         14, True, BLACK, WARN_FILL, None, space=1.5)


# --------------------------------------------------- 7 第5条 JIRA/Winchill ----
def systems(prs, n):
    sl = page(prs, "[第5条] JIRAおよびWinchillの役割分担",
              "JIRAでは版の履歴を都度確認できないため、バージョン管理の観点から\n"
              "Winchillにて一元管理する。", (DECIDED,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表2　システム別の保管対象")
    cols = [("システム", 1.70), ("保管対象", 4.10), ("用途", 3.10), ("アクセス", W - 8.90)]
    thead(sl, L, y + 0.38, cols)
    rows = [("JIRA", "決定事項／進捗報告／議事録／\nステータス／要件確認",
             "プロジェクト管理\nコミュニケーション記録", "生準メンバ\n（今後、関連部署）"),
            ("Winchill", "図面・設計図／CADファイル／技術仕様書／\n製造情報／その他機密資料",
             "機密情報管理\n設計・技術情報の版管理", "権限者のみ\n（別途制御）")]
    ry = y + 0.74
    for name, keep, use, acc in rows:
        tcell(sl, L, ry, 1.70, 1.20, name, 14, True, DARK, PALE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.70, ry, 4.10, 1.20, keep, 11.5, space=1.35)
        tcell(sl, L + 5.80, ry, 3.10, 1.20, use, 11.5, space=1.35)
        tcell(sl, L + 8.90, ry, W - 8.90, 1.20, acc, 11.5, space=1.35)
        ry += 1.20
    note(sl, L, ry + 0.28, W, 1.00, "図面・CADファイルの取扱い",
         "JIRAには登録せず、Winchillにて保管する。JIRAには参照リンクまたは"
         "Winchill No.のみを記載する。JIRAでは版の履歴を都度確認できないため、"
         "バージョン管理はWinchillに一元化する。")


# ------------------------------------------------------- 8 第6条 情報区分 ----
def classes(prs, n):
    sl = page(prs, "[第6条] 情報区分と取扱い",
              "機密区分の情報を参照する必要がある場合は、JIRAには外部リンクのみを記載する。\n"
              "添付ファイルとしては登録しない。", (DECIDED,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "表3　情報区分別の登録可否")
    cols = [("情報区分", 3.30), ("JIRAへの登録", 2.40), ("保管先", W - 5.70)]
    thead(sl, L, y + 0.38, cols)
    rows = [("一般（Public）", "可", GREEN, "JIRA"),
            ("社内（Internal）", "可", GREEN, "JIRA（アクセス権限を設定する）"),
            ("機密（Confidential）", "不可", RED, "Winchill／専用ストレージ　※参照リンクのみ"),
            ("極秘（Secret）", "不可", RED, "限定者のみアクセス可能な場所")]
    ry = y + 0.74
    for k, ok, col, dest in rows:
        tcell(sl, L, ry, 3.30, 0.64, k, 12.5, True, DARK, PALE, pad=0.16)
        tcell(sl, L + 3.30, ry, 2.40, 0.64, ok, 14, True, col, WHITE, PP_ALIGN.CENTER)
        tcell(sl, L + 5.70, ry, W - 5.70, 0.64, dest, 12, pad=0.16)
        ry += 0.64
    note(sl, L, ry + 0.30, W, 0.88, "運用",
         "機密区分の情報を参照する必要がある場合は、JIRAには外部リンクのみを記載し、"
         "添付ファイルとしては登録しない。")


# ------------------------------------------------- 9 第7条 チェックリスト ----
def checklist(prs, n):
    sl = page(prs, "[第7条] 登録前チェックリスト",
              "JIRAに文書を登録する前に、以下の5点を確認する。", (DECIDED,), n)
    y = BODY_TOP
    secbar(sl, L, y, W, "登録前の確認事項")
    items = [("図面・CADファイルではないか", "Winchillで保管すべきものではないか"),
             ("顧客の個人情報・連絡先が含まれていないか", "個人携帯番号・自宅住所・私的な内容"),
             ("必要なアクセス権限が設定されているか", "社内区分は権限設定のうえ登録する"),
             ("ファイルサイズは50MB以下か", "超過する場合は分割または別保管を検討する"),
             ("関連部署に閲覧されても差し支えない内容か", "公開範囲の拡大を前提に判断する")]
    ry = y + 0.38
    for i, (q, rem) in enumerate(items, 1):
        tcell(sl, L, ry, 0.56, 0.64, "☐", 15, align=PP_ALIGN.CENTER)
        tcell(sl, L + 0.56, ry, 0.56, 0.64, str(i), 12, True, WHITE, BLUE, PP_ALIGN.CENTER)
        tcell(sl, L + 1.12, ry, 5.98, 0.64, q, 12.5, True, DARK, PALE, pad=0.16)
        tcell(sl, L + 7.10, ry, W - 7.10, 0.64, rem, 11.5, pad=0.16)
        ry += 0.64
    note(sl, L, ry + 0.30, W, 0.80, "図面・設計図の場合",
         "Winchillに登録のうえ、JIRAには参照リンクまたはWinchill No.を記載する。",
         fill=WARN_FILL)


# ---------------------------------------------------------------- 10 まとめ ----
def closing(prs, n):
    sl = page(prs, "まとめ",
              "本規程の要点と、制定にあたって残っている未決事項。", (DECIDED, OPEN), n)
    y = BODY_TOP
    cell(sl, L, y, W, 0.80,
         "記録の所在を一つに定めることが、認識の相違と属人化を防ぐ",
         22, True, WHITE, BLUE, None)
    pts = [("残す場所", "個人のメールボックスではなく、\n当該タスクのチケットに保管する。"),
           ("残し方", "社外メールは原本のまま。口頭決定は\n確認メールを送り、それを保管する。"),
           ("分けるもの", "図面・CADと機密以上はWinchillへ。\nJIRAには参照リンクのみを記載する。")]
    cw = (W - 0.40) / 3
    for i, (t, d) in enumerate(pts):
        x = L + i * (cw + 0.20)
        cell(sl, x, y + 1.08, cw, 0.40, t, 13, True, WHITE, DARK, None)
        tcell(sl, x, y + 1.48, cw, 1.10, d, 12, align=PP_ALIGN.CENTER, space=1.35)
    ny = y + 2.90
    secbar(sl, L, ny, W, "未決事項")
    opens = [("01", "文書番号の採番"), ("02", "管理部署の確定"),
             ("03", "関連部署への公開時期（別途通知）")]
    ry = ny + 0.38
    for no, t in opens:
        tcell(sl, L, ry, 0.66, 0.44, no, 11.5, True, BLACK, TAG_OPEN, PP_ALIGN.CENTER)
        tcell(sl, L + 0.66, ry, W - 0.66, 0.44, t, 12, pad=0.16)
        ry += 0.44
    say(sl, L, ry + 0.14, W, 0.28,
        "附則　本規程は制定日より実施する。改定は管理部署の承認を経て行う。",
        12, False, MUTED)


def back_cover(prs):
    """Back Cover layout. Its title placeholder carries an unusual idx that
    python-pptx does not clone, so the wordmark is drawn at the layout's own
    measured position (x=3.82 y=2.73 w=5.68 h=2.61)."""
    sl = prs.slides.add_slide(prs.slide_master.slide_layouts[5])
    say(sl, 3.82, 2.73, 5.68, 2.61, "Thank you", 32, True, WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return sl


def build(template, output):
    prs = Presentation(base_from_template(template))
    cover(prs)
    for i, fn in enumerate([overview, purpose, scope, targets, external,
                            systems, classes, checklist, closing], start=2):
        fn(prs, i)
    back_cover(prs)
    prs.save(output)
    print("saved", output, "-", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("-o", "--output", default="jira_document_policy_daicel_tpl.pptx")
    a = ap.parse_args()
    build(a.template, a.output)
