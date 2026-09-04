# -*- coding: utf-8 -*-
"""Page system measured from the Bジ／TITUS delivery deck (Accenture Japan style).

Measured from the source deck, not invented:
  資料区分   18pt bold Meiryo UI   x=0.67 y=0.11 w=9.69 h=0.38 (bottom-anchored)
  タイトル   24pt bold Meiryo UI   x=0.67 y=0.52 w=12.00 h=0.47
  説明文     18pt Meiryo UI        x=0.67 y=0.99 w=11.51 h=0.61 (max 2 lines)
  プロセスID 16pt Meiryo UI        x=10.76 y=0.04 w=2.52 h=0.37
  ステータスタグ 12pt, w=1.34 h=0.31, black 1pt border, base colour at lumMod20/lumOff80
              rightmost x=11.63, next x=10.29
  セクションバー accent4 000088, white 18pt centred, black 0.75pt border, h=0.39
  フッター   Copyright © 2021 Accenture...  x=7.58 y=7.24  /  page no x=12.70 y=7.22
Theme: accent4 000088 (navy), accent2 A100FF, dk2 FF9128; body font Meiryo UI.
"""
import colorsys

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def C(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def tint(hexstr, lum_mod=0.20, lum_off=0.80):
    """Reproduce OOXML <a:lumMod>/<a:lumOff> on an sRGB colour."""
    r, g, b = (int(hexstr[i:i + 2], 16) / 255 for i in (0, 2, 4))
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    l = min(1.0, l * lum_mod + lum_off)
    r, g, b = colorsys.hls_to_rgb(h, l, s)
    return RGBColor(round(r * 255), round(g * 255), round(b * 255))


BLACK, WHITE = C("000000"), C("FFFFFF")
NAVY = C("000088")        # accent4 - section bars, table headers
BLUE = C("0064D2")        # key-message bar
PALE = C("E7EAF7")        # table sub-header / label cells
GREY = C("E5E5E5")
MIDGREY = C("919191")     # lt2
RED = C("FF0000")
GREEN = C("009900")
ORANGE = "FA9628"         # dk2, tag base
PURPLE = "644080"         # tag / annotation base
TAG_DECIDED = tint(ORANGE)
TAG_OPEN = tint(PURPLE)
NOTE_FILL = tint(PURPLE)
WARN_FILL = tint(ORANGE)

FONT = "Meiryo UI"

# page grid
L, R = 0.67, 12.67
W = R - L
BODY_TOP = 1.70
FOOT = "Copyright © 2021 Accenture. All rights reserved. Highly Confidential"


def style(run, size, bold=False, color=BLACK, face=FONT):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, face
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", face)


def textbox(sl, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tb, tf


def say(sl, x, y, w, h, text, size, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
    _, tf = textbox(sl, x, y, w, h, align, anchor)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.line_spacing = space
        r = p.add_run()
        r.text = ln
        style(r, size, bold, color)
    return tf


def rich(sl, x, y, w, h, parts, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=None):
    _, tf = textbox(sl, x, y, w, h, align, anchor)
    p = tf.paragraphs[0]
    if space:
        p.line_spacing = space
    for text, size, bold, color in parts:
        r = p.add_run()
        r.text = text
        style(r, size, bold, color)
    return tf


def line(sl, x1, y1, x2, y2, color, width=1.0, dash=None):
    ln = sl.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        el = ln.line._get_or_add_ln()
        el.append(el.makeelement(qn("a:prstDash"), {"val": dash}))
    return ln


def cell(sl, x, y, w, h, text="", size=12, bold=False, fg=BLACK, fill=WHITE,
         border=BLACK, bw=0.75, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
         space=None, pad=0.10):
    """A bordered box — the deck's basic table/diagram unit."""
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border
        sh.line.width = Pt(bw)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space:
            p.line_spacing = space
        if ln:
            r = p.add_run()
            r.text = ln
            style(r, size, bold, fg)
    return sh


def secbar(sl, x, y, w, text, h=0.39, size=18, fill=NAVY, fg=WHITE):
    """Navy section bar: accent4 fill, white centred label, hairline border."""
    return cell(sl, x, y, w, h, text, size, False, fg, fill, BLACK, 0.75)


def note(sl, x, y, w, h, label, body, fill=NOTE_FILL, size=13, lsize=13):
    """Annotation block — the deck's pale tinted callout with a bold lead-in."""
    cell(sl, x, y, w, h, "", fill=fill, border=BLACK, bw=0.75)
    parts = ([(label + "　", lsize, True, BLACK)] if label else []) + \
            [(body, size, False, BLACK)]
    rich(sl, x + 0.16, y + 0.11, w - 0.32, h - 0.22, parts,
         anchor=MSO_ANCHOR.MIDDLE, space=1.3)


def frame(sl, doc_class, title, desc=None, ref=None, tags=(), page_no=None):
    """The standard page furniture, at the measured positions."""
    say(sl, L, 0.11, 9.69, 0.38, doc_class, 18, True, BLACK, anchor=MSO_ANCHOR.BOTTOM)
    say(sl, L, 0.52, 12.00, 0.47, title, 24, True, BLACK)
    if desc:
        say(sl, L, 0.99, 11.51, 0.61, desc, 16, False, BLACK, space=1.25)
    if ref:
        say(sl, 10.76, 0.04, 2.52, 0.37, ref, 16, False, BLACK, align=PP_ALIGN.RIGHT)
    for i, (label, fill) in enumerate(reversed(list(tags))):
        cell(sl, 11.63 - i * 1.34, 0.44, 1.34, 0.31, label, 12, False, BLACK,
             fill, BLACK, 1.0)
    say(sl, 7.58, 7.24, 4.81, 0.20, FOOT, 9, False, BLACK, align=PP_ALIGN.RIGHT)
    if page_no is not None:
        say(sl, 12.70, 7.22, 0.30, 0.20, str(page_no), 12, False, BLACK)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])
