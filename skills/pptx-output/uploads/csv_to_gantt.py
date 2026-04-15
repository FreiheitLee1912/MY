"""
CSV to Gantt Chart (PPTX) Generator
===================================
Reads Jira CSV exports and generates PPT-ready Gantt charts.
Output files are saved to the ./output/ folder.

Usage:
    python csv_to_gantt.py                          # Process all CSV files in Downloads
    python csv_to_gantt.py path/to/file.csv         # Process a specific CSV file
    python csv_to_gantt.py file1.csv file2.csv      # Process multiple CSV files
"""

import csv
import os
import sys
import re
import io
from datetime import datetime, timedelta
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Fix Windows console encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# ========================================
# Configuration
# ========================================
SCRIPT_DIR = Path(__file__).parent
OUTPUT_DIR = SCRIPT_DIR / "output"
DOWNLOADS_DIR = Path.home() / "Downloads"

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)

# Layout constants
LEFT_MARGIN = Inches(0.3)
TOP_MARGIN = Inches(0.6)
HEADER_HEIGHT = Inches(0.3)
ROW_HEIGHT = Inches(0.3)
LABEL_WIDTH = Inches(3.2)
CHART_LEFT = LEFT_MARGIN + LABEL_WIDTH + Inches(0.1)
CHART_WIDTH = SLIDE_WIDTH - CHART_LEFT - Inches(0.3)
TASKS_PER_SLIDE = 18

# Colors (RGB)
BG_COLOR = RGBColor(0x0f, 0x11, 0x17)
BG_SECONDARY = RGBColor(0x16, 0x18, 0x22)
BG_ROW_ALT = RGBColor(0x1a, 0x1c, 0x28)
TEXT_PRIMARY = RGBColor(0xe8, 0xea, 0xed)
TEXT_SECONDARY = RGBColor(0x9a, 0xa0, 0xa6)
TEXT_MUTED = RGBColor(0x5f, 0x63, 0x68)
GRID_COLOR = RGBColor(0x2a, 0x2d, 0x3a)
TODAY_COLOR = RGBColor(0xeb, 0x57, 0x57)

# Task type colors
TYPE_COLORS = {
    'Milestone':   {'rgb': RGBColor(0xf2, 0x99, 0x4a), 'label': 'MS'},
    'Goal':        {'rgb': RGBColor(0x6f, 0xcf, 0x97), 'label': 'Goal'},
    'Review':      {'rgb': RGBColor(0x56, 0xcc, 0xf2), 'label': 'Rev'},
    'Event':       {'rgb': RGBColor(0xbb, 0x6b, 0xd9), 'label': 'Evt'},
    'タスク':       {'rgb': RGBColor(0x8a, 0xb4, 0xf8), 'label': 'Task'},
    'Task':        {'rgb': RGBColor(0x8a, 0xb4, 0xf8), 'label': 'Task'},
    'Sub task':    {'rgb': RGBColor(0x8a, 0xb4, 0xf8), 'label': 'Sub'},
    'Validation':  {'rgb': RGBColor(0xeb, 0x57, 0x57), 'label': 'Val'},
    'SOP':         {'rgb': RGBColor(0xf2, 0xc9, 0x4c), 'label': 'SOP'},
    'Certificate': {'rgb': RGBColor(0x2f, 0x80, 0xed), 'label': 'Cert'},
    'PPAP':        {'rgb': RGBColor(0xff, 0x6b, 0x6b), 'label': 'PPAP'},
    'Top-level initiative': {'rgb': RGBColor(0x8a, 0xb4, 0xf8), 'label': 'Init'},
}

DEFAULT_TYPE = {'rgb': RGBColor(0x82, 0x82, 0x82), 'label': '?'}
COMPLETED_COLOR = RGBColor(0x27, 0xae, 0x60)


# ========================================
# Date Parsing
# ========================================
def parse_date(date_str):
    """Parse various date formats from Jira CSV exports."""
    if not date_str or not date_str.strip():
        return None
    
    s = date_str.strip()
    # Remove time portions: "12:00 午前", "12:00 AM", etc.
    s = re.sub(r'\s+\d{1,2}:\d{2}(:\d{2})?\s*(午前|午後|AM|PM)?', '', s, flags=re.IGNORECASE).strip()
    
    # Format: DD/M/YY  (Jira: day first, 2-digit year last, e.g. "20/4/26" = Apr 20 2026)
    m = re.match(r'^(\d{1,4})/(\d{1,2})/(\d{1,4})$', s)
    if m:
        a, b, c = int(m.group(1)), int(m.group(2)), int(m.group(3))
        if a > 31:
            # YYYY/M/D
            year, month, day = a, b, c
        elif c <= 99:
            # DD/M/YY — Jira format
            day, month, year = a, b, c + 2000
        else:
            # D/M/YYYY
            day, month, year = a, b, c
        try:
            return datetime(year, month, day)
        except ValueError:
            return None
    
    # Try YYYY-MM-DD
    m = re.match(r'^(\d{4})-(\d{1,2})-(\d{1,2})$', s)
    if m:
        try:
            return datetime(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            return None
    
    return None



# ========================================
# CSV Reading
# ========================================
def read_csv(filepath):
    """Read a Jira CSV export and return a list of task dictionaries."""
    tasks = []
    filepath = Path(filepath)
    
    # Try different encodings
    for encoding in ['utf-8', 'utf-8-sig', 'cp932', 'shift_jis', 'latin-1']:
        try:
            with open(filepath, 'r', encoding=encoding) as f:
                reader = csv.reader(f)
                rows = list(reader)
            # Verify we got valid data
            if rows and len(rows[0]) >= 4:
                break
        except (UnicodeDecodeError, UnicodeError):
            continue
    else:
        print(f"  ⚠ Could not read {filepath.name} with any encoding")
        return []
    
    if not rows:
        return []
    
    # First row is header
    headers = rows[0]
    
    for row in rows[1:]:
        if len(row) < 4:
            continue
        
        key = row[0].strip() if len(row) > 0 else ''
        task_type = row[1].strip() if len(row) > 1 else ''
        parent = row[2].strip() if len(row) > 2 else ''
        summary = row[3].strip() if len(row) > 3 else ''
        status = row[4].strip() if len(row) > 4 else 'To Do'
        assignee = row[5].strip() if len(row) > 5 else ''
        start_date_str = row[6].strip() if len(row) > 6 else ''
        est_start_str = row[7].strip() if len(row) > 7 else ''
        deadline_str = row[8].strip() if len(row) > 8 else ''
        est_deadline_str = row[9].strip() if len(row) > 9 else ''
        
        if not key:
            continue
        
        start_date = parse_date(start_date_str) or parse_date(est_start_str)
        end_date = parse_date(deadline_str) or parse_date(est_deadline_str)
        
        tasks.append({
            'key': key,
            'type': task_type or 'Task',
            'parent': parent,
            'summary': summary or key,
            'status': status,
            'assignee': assignee,
            'start_date': start_date,
            'end_date': end_date,
            'project': key.split('-')[0],
        })
    
    return tasks


# ========================================
# Task Grouping
# ========================================
def group_tasks_by_type(tasks):
    """Group tasks by type, sorted by type order then by start date."""
    type_order = [
        'Milestone', 'Top-level initiative', 'Goal', 'Review', 'Event',
        'SOP', 'Certificate', 'Validation', 'PPAP', 'タスク', 'Task', 'Sub task'
    ]
    
    # Only tasks with dates
    with_dates = [t for t in tasks if t['start_date'] or t['end_date']]
    
    groups = {}
    for task in with_dates:
        g = task['type']
        if g not in groups:
            groups[g] = []
        groups[g].append(task)
    
    # Sort groups by type order
    def type_sort_key(name):
        try:
            return type_order.index(name)
        except ValueError:
            return 999
    
    result = []
    for group_name in sorted(groups.keys(), key=type_sort_key):
        group_tasks = sorted(groups[group_name], key=lambda t: t['start_date'] or datetime(9999, 1, 1))
        result.append({'is_group': True, 'label': group_name, 'count': len(group_tasks)})
        result.extend(group_tasks)
    
    return result


# ========================================
# Timeline Calculations
# ========================================
def get_timeline_bounds(tasks):
    """Get the earliest start and latest end dates from tasks."""
    dates = []
    for t in tasks:
        if isinstance(t, dict) and not t.get('is_group'):
            if t.get('start_date'):
                dates.append(t['start_date'])
            if t.get('end_date'):
                dates.append(t['end_date'])
    
    if not dates:
        now = datetime.now()
        return now, now + timedelta(days=365)
    
    earliest = min(dates)
    latest = max(dates)
    
    # Padding: 1 month before and after
    start = datetime(earliest.year, earliest.month, 1) - timedelta(days=30)
    start = datetime(start.year, start.month, 1)
    
    if latest.month == 12:
        end = datetime(latest.year + 1, 2, 1)
    else:
        end_month = latest.month + 2
        end_year = latest.year
        if end_month > 12:
            end_month -= 12
            end_year += 1
        end = datetime(end_year, end_month, 1)
    
    return start, end


def get_months_between(start, end):
    """Get list of (year, month) tuples between start and end."""
    months = []
    current = datetime(start.year, start.month, 1)
    while current < end:
        months.append(current)
        if current.month == 12:
            current = datetime(current.year + 1, 1, 1)
        else:
            current = datetime(current.year, current.month + 1, 1)
    return months


def days_in_month(dt):
    """Get number of days in a month."""
    if dt.month == 12:
        next_month = datetime(dt.year + 1, 1, 1)
    else:
        next_month = datetime(dt.year, dt.month + 1, 1)
    return (next_month - datetime(dt.year, dt.month, 1)).days


# ========================================
# PPTX Generation
# ========================================
def make_transparent(color, alpha=0.15):
    """Create a lighter version of a color for backgrounds."""
    r = int(color[0] + (255 - color[0]) * (1 - alpha))
    g = int(color[1] + (255 - color[1]) * (1 - alpha))
    b = int(color[2] + (255 - color[2]) * (1 - alpha))
    return RGBColor(r, g, b)


def add_text_box(slide, text, x, y, w, h, 
                 font_size=7, color=TEXT_PRIMARY, bold=False, 
                 align=PP_ALIGN.LEFT, font_name='Meiryo UI'):
    """Helper to add a text box with common formatting."""
    from pptx.util import Emu as _Emu
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Vertical alignment
    tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    # Vertical center
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.bodyPr
    bodyPr.set('anchor', 'ctr')
    return txBox


def generate_pptx(tasks, output_path, title="Gantt Chart"):
    """Generate a PPTX file with Gantt chart from tasks."""
    
    grouped = group_tasks_by_type(tasks)
    if not grouped:
        print("  ⚠ No tasks with dates found.")
        return False
    
    # Get only non-group items for timeline
    task_items = [t for t in grouped if not t.get('is_group')]
    bounds_start, bounds_end = get_timeline_bounds(task_items)
    total_days = (bounds_end - bounds_start).days
    if total_days <= 0:
        total_days = 1
    
    months = get_months_between(bounds_start, bounds_end)
    today = datetime.now()
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Paginate tasks (keeping group headers)
    pages = []
    current_page = []
    task_count = 0
    
    for item in grouped:
        if item.get('is_group'):
            current_page.append(item)
        else:
            current_page.append(item)
            task_count += 1
            if task_count >= TASKS_PER_SLIDE:
                pages.append(current_page)
                current_page = []
                task_count = 0
    
    if current_page:
        pages.append(current_page)
    
    if not pages:
        print("  ⚠ No data to render.")
        return False
    
    for page_idx, page_items in enumerate(pages):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Title
        page_label = f" ({page_idx + 1}/{len(pages)})" if len(pages) > 1 else ""
        add_text_box(slide, f"{title}{page_label}",
                     LEFT_MARGIN, Inches(0.15), Inches(10), Inches(0.35),
                     font_size=14, color=TEXT_PRIMARY, bold=True)
        
        # ---- Month Headers ----
        chart_width_emu = CHART_WIDTH
        month_x = CHART_LEFT
        
        for month_dt in months:
            days = days_in_month(month_dt)
            month_frac = days / total_days
            month_width = int(chart_width_emu * month_frac)
            
            # Month header box
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, month_x, TOP_MARGIN, month_width, HEADER_HEIGHT
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = BG_SECONDARY
            shape.line.color.rgb = GRID_COLOR
            shape.line.width = Pt(0.5)
            
            # Month label
            tf = shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            month_names = ['1月','2月','3月','4月','5月','6月','7月','8月','9月','10月','11月','12月']
            p.text = f"{month_dt.year}/{month_names[month_dt.month - 1]}"
            p.font.size = Pt(7)
            p.font.color.rgb = TEXT_SECONDARY
            p.font.name = 'Meiryo UI'
            p.alignment = PP_ALIGN.CENTER
            from pptx.oxml.ns import qn
            tf._txBody.bodyPr.set('anchor', 'ctr')
            
            # Vertical grid line
            grid_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, month_x, TOP_MARGIN + HEADER_HEIGHT,
                Pt(1), Inches(len(page_items) * 0.3 + 0.1)
            )
            grid_line.fill.solid()
            grid_line.fill.fore_color.rgb = GRID_COLOR
            grid_line.line.fill.background()
            
            month_x += month_width
        
        # ---- Today Line ----
        today_days = (today - bounds_start).days
        if 0 <= today_days <= total_days:
            today_frac = today_days / total_days
            today_x = int(CHART_LEFT + chart_width_emu * today_frac)
            
            today_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, today_x, TOP_MARGIN - Inches(0.05),
                Pt(2), Inches(len(page_items) * 0.3 + 0.5)
            )
            today_line.fill.solid()
            today_line.fill.fore_color.rgb = TODAY_COLOR
            today_line.line.fill.background()
            
            add_text_box(slide, 'Today',
                         today_x - Inches(0.2), TOP_MARGIN - Inches(0.2),
                         Inches(0.4), Inches(0.15),
                         font_size=5, color=TODAY_COLOR, align=PP_ALIGN.CENTER)
        
        # ---- Task Rows ----
        row_y = TOP_MARGIN + HEADER_HEIGHT + Inches(0.05)
        row_idx = 0
        
        for item in page_items:
            y = row_y + int(row_idx * ROW_HEIGHT)
            
            if item.get('is_group'):
                # Group header row
                bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, LEFT_MARGIN, y,
                    SLIDE_WIDTH - LEFT_MARGIN * 2, ROW_HEIGHT
                )
                bg.fill.solid()
                bg.fill.fore_color.rgb = RGBColor(0x1e, 0x22, 0x33)
                bg.line.fill.background()
                
                type_conf = TYPE_COLORS.get(item['label'], DEFAULT_TYPE)
                add_text_box(slide, f"▸ {item['label']} ({item['count']})",
                             LEFT_MARGIN + Inches(0.1), y,
                             LABEL_WIDTH, ROW_HEIGHT,
                             font_size=7, color=type_conf['rgb'], bold=True)
                row_idx += 1
                continue
            
            # Alternating row background
            if row_idx % 2 == 0:
                bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, LEFT_MARGIN, y,
                    SLIDE_WIDTH - LEFT_MARGIN * 2, ROW_HEIGHT
                )
                bg.fill.solid()
                bg.fill.fore_color.rgb = BG_ROW_ALT
                bg.line.fill.background()
            
            type_conf = TYPE_COLORS.get(item['type'], DEFAULT_TYPE)
            
            # Task key
            add_text_box(slide, item['key'],
                         LEFT_MARGIN + Inches(0.05), y,
                         Inches(0.7), ROW_HEIGHT,
                         font_size=6, color=TEXT_MUTED)
            
            # Type badge
            badge_w = Inches(0.4)
            badge_h = Inches(0.18)
            badge_y = y + (ROW_HEIGHT - badge_h) // 2
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                LEFT_MARGIN + Inches(0.75), badge_y,
                badge_w, badge_h
            )
            badge.fill.solid()
            # Lighter version of type color
            r, g, b = type_conf['rgb'][0], type_conf['rgb'][1], type_conf['rgb'][2]
            badge.fill.fore_color.rgb = RGBColor(
                min(255, int(r * 0.2 + BG_COLOR[0] * 0.8)),
                min(255, int(g * 0.2 + BG_COLOR[1] * 0.8)),
                min(255, int(b * 0.2 + BG_COLOR[2] * 0.8))
            )
            badge.line.fill.background()
            tf = badge.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = type_conf['label']
            p.font.size = Pt(5)
            p.font.color.rgb = type_conf['rgb']
            p.font.bold = True
            p.font.name = 'Meiryo UI'
            p.alignment = PP_ALIGN.CENTER
            from pptx.oxml.ns import qn
            tf._txBody.bodyPr.set('anchor', 'ctr')
            
            # Task summary
            add_text_box(slide, item['summary'],
                         LEFT_MARGIN + Inches(1.2), y,
                         LABEL_WIDTH - Inches(1.2), ROW_HEIGHT,
                         font_size=7, color=TEXT_PRIMARY)
            
            # ---- Draw bar ----
            if item['start_date'] or item['end_date']:
                start = item['start_date'] or item['end_date']
                end = item['end_date'] or item['start_date']
                
                start_days = (start - bounds_start).days
                duration_days = max((end - start).days, 1)
                
                bar_x = int(CHART_LEFT + chart_width_emu * (start_days / total_days))
                bar_w = max(int(chart_width_emu * (duration_days / total_days)), Inches(0.06))
                bar_h = Inches(0.18)
                bar_y = y + (ROW_HEIGHT - bar_h) // 2
                
                is_milestone = item['type'] == 'Milestone'
                is_completed = item['status'] in ('完了', 'Done')
                bar_color = COMPLETED_COLOR if is_completed else type_conf['rgb']
                
                if is_milestone:
                    # Diamond for milestones
                    diamond_size = Inches(0.15)
                    diamond = slide.shapes.add_shape(
                        MSO_SHAPE.DIAMOND,
                        bar_x - diamond_size // 2, bar_y,
                        diamond_size, bar_h
                    )
                    diamond.fill.solid()
                    diamond.fill.fore_color.rgb = bar_color
                    diamond.line.fill.background()
                else:
                    # Regular bar
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        bar_x, bar_y, bar_w, bar_h
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = bar_color
                    bar.line.fill.background()
                    
                    # Bar label (if wide enough)
                    if bar_w > Inches(0.8):
                        tf = bar.text_frame
                        tf.word_wrap = False
                        p = tf.paragraphs[0]
                        p.text = item['summary']
                        p.font.size = Pt(5)
                        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                        p.font.name = 'Meiryo UI'
                        p.alignment = PP_ALIGN.LEFT
                        tf._txBody.bodyPr.set('anchor', 'ctr')
                    
                    # Date labels on bar sides
                    if start_days >= 0:
                        date_str = start.strftime('%m/%d')
                        add_text_box(slide, date_str,
                                     bar_x - Inches(0.35), bar_y,
                                     Inches(0.33), bar_h,
                                     font_size=4, color=TEXT_MUTED,
                                     align=PP_ALIGN.RIGHT)
                    
                    if duration_days > 1:
                        date_str = end.strftime('%m/%d')
                        add_text_box(slide, date_str,
                                     bar_x + bar_w + Inches(0.02), bar_y,
                                     Inches(0.33), bar_h,
                                     font_size=4, color=TEXT_MUTED,
                                     align=PP_ALIGN.LEFT)
            
            row_idx += 1
        
        # ---- Legend ----
        legend_y = SLIDE_HEIGHT - Inches(0.35)
        legend_x = LEFT_MARGIN
        
        legend_items = [
            ('Milestone', TYPE_COLORS['Milestone']['rgb']),
            ('Goal', TYPE_COLORS['Goal']['rgb']),
            ('Review', TYPE_COLORS['Review']['rgb']),
            ('Event', TYPE_COLORS['Event']['rgb']),
            ('Task', TYPE_COLORS['Task']['rgb']),
            ('Validation', TYPE_COLORS['Validation']['rgb']),
            ('Completed', COMPLETED_COLOR),
        ]
        
        for label, color in legend_items:
            # Color swatch
            swatch = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                legend_x, legend_y + Inches(0.03),
                Inches(0.2), Inches(0.1)
            )
            swatch.fill.solid()
            swatch.fill.fore_color.rgb = color
            swatch.line.fill.background()
            
            # Label
            add_text_box(slide, label,
                         legend_x + Inches(0.22), legend_y,
                         Inches(0.6), Inches(0.15),
                         font_size=6, color=TEXT_SECONDARY)
            
            legend_x += Inches(0.9)
        
        # Today legend
        today_swatch = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            legend_x, legend_y,
            Pt(2), Inches(0.15)
        )
        today_swatch.fill.solid()
        today_swatch.fill.fore_color.rgb = TODAY_COLOR
        today_swatch.line.fill.background()
        
        add_text_box(slide, 'Today',
                     legend_x + Inches(0.05), legend_y,
                     Inches(0.5), Inches(0.15),
                     font_size=6, color=TODAY_COLOR)
    
    # Save
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return True


# ========================================
# Main
# ========================================
def find_csv_files():
    """Find relevant CSV files in Downloads folder."""
    csv_files = []
    if DOWNLOADS_DIR.exists():
        for f in DOWNLOADS_DIR.glob("*.csv"):
            # Look for Jira-style exports and specific project CSVs
            name_lower = f.name.lower()
            if any(kw in name_lower for kw in ['jira', 'dpft', 'dat', 'dss', 'dssc', 'tbd']):
                csv_files.append(f)
    return sorted(csv_files, key=lambda f: f.stat().st_mtime, reverse=True)


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    print("=" * 60)
    print("  📊 CSV → Gantt Chart (PPTX) Generator")
    print("=" * 60)
    
    # Determine input files
    if len(sys.argv) > 1:
        csv_files = [Path(f) for f in sys.argv[1:] if Path(f).exists()]
    else:
        csv_files = find_csv_files()
    
    if not csv_files:
        print("\n  ⚠ No CSV files found!")
        print(f"  Looking in: {DOWNLOADS_DIR}")
        print("  Usage: python csv_to_gantt.py [file1.csv] [file2.csv]")
        return
    
    print(f"\n  Found {len(csv_files)} CSV file(s):")
    for f in csv_files:
        print(f"    📄 {f.name}")
    
    # Process each file individually
    for csv_file in csv_files:
        print(f"\n  Processing: {csv_file.name}")
        tasks = read_csv(csv_file)
        
        if not tasks:
            print(f"    ⚠ No valid tasks found in {csv_file.name}")
            continue
        
        # Count tasks with dates
        with_dates = sum(1 for t in tasks if t['start_date'] or t['end_date'])
        print(f"    ✓ {len(tasks)} tasks loaded ({with_dates} with dates)")
        
        # Get the project prefix for the title
        projects = set(t['project'] for t in tasks)
        title = " / ".join(sorted(projects)) + " Gantt Chart"
        
        # Generate output filename
        stem = csv_file.stem.replace(' ', '_')
        output_path = OUTPUT_DIR / f"{stem}_gantt.pptx"
        
        success = generate_pptx(tasks, output_path, title=title)
        
        if success:
            print(f"    ✅ Saved: {output_path}")
        else:
            print(f"    ❌ Failed to generate PPTX")
    
    # Also generate a combined file if multiple CSVs
    if len(csv_files) > 1:
        print(f"\n  Generating combined Gantt chart...")
        all_tasks = []
        for csv_file in csv_files:
            all_tasks.extend(read_csv(csv_file))
        
        if all_tasks:
            projects = set(t['project'] for t in all_tasks)
            combined_title = " + ".join(sorted(projects)) + " Combined Gantt Chart"
            combined_path = OUTPUT_DIR / "combined_gantt.pptx"
            
            success = generate_pptx(all_tasks, combined_path, title=combined_title)
            if success:
                print(f"    ✅ Saved: {combined_path}")
    
    print(f"\n  📁 Output folder: {OUTPUT_DIR.resolve()}")
    print("=" * 60)


if __name__ == '__main__':
    main()
