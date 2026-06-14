"""
Custom HTTP Server for Gantt Chart Generator
=============================================
Serves static files and handles PPTX generation via POST /api/generate-pptx
Output files are saved to ./output/ folder.

Usage:
    python server.py
    Open http://localhost:8090
"""

import http.server
import json
import os
import sys
import io
import re
import csv
from datetime import datetime, timedelta
from pathlib import Path
from urllib.parse import unquote

# Fix Windows console encoding
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = Path(__file__).parent
OUTPUT_DIR = SCRIPT_DIR / "output"
PORT = 8090

# ========================================
# PPTX Color Config (Light Theme for PPT)
# ========================================
# Light theme colors - white/light background, dark text
BG_COLOR = RGBColor(0xff, 0xff, 0xff)         # White background
BG_SECONDARY = RGBColor(0xf0, 0xf2, 0xf5)     # Light gray header
BG_ROW_ALT = RGBColor(0xf8, 0xf9, 0xfb)       # Very light row alt
TEXT_PRIMARY = RGBColor(0x1a, 0x1a, 0x2e)       # Dark text
TEXT_SECONDARY = RGBColor(0x64, 0x74, 0x8b)     # Medium gray
TEXT_MUTED = RGBColor(0x94, 0xa3, 0xb8)         # Light gray
GRID_COLOR = RGBColor(0xe2, 0xe8, 0xf0)         # Light grid
TODAY_COLOR = RGBColor(0xeb, 0x57, 0x57)         # Red today line
TITLE_COLOR = RGBColor(0x1e, 0x3a, 0x5f)        # Dark blue title

TYPE_COLORS = {
    'Milestone':   {'rgb': RGBColor(0xf2, 0x99, 0x4a), 'label': 'MS'},
    'Goal':        {'rgb': RGBColor(0x4c, 0xaf, 0x50), 'label': 'Goal'},
    'Review':      {'rgb': RGBColor(0x29, 0x96, 0xd6), 'label': 'Rev'},
    'Event':       {'rgb': RGBColor(0x9c, 0x27, 0xb0), 'label': 'Evt'},
    'タスク':       {'rgb': RGBColor(0x42, 0x7a, 0xd6), 'label': 'Task'},
    'Task':        {'rgb': RGBColor(0x42, 0x7a, 0xd6), 'label': 'Task'},
    'Sub task':    {'rgb': RGBColor(0x64, 0x95, 0xed), 'label': 'Sub'},
    'Validation':  {'rgb': RGBColor(0xe5, 0x3e, 0x3e), 'label': 'Val'},
    'SOP':         {'rgb': RGBColor(0xe6, 0xa8, 0x17), 'label': 'SOP'},
    'Certificate': {'rgb': RGBColor(0x2f, 0x80, 0xed), 'label': 'Cert'},
    'PPAP':        {'rgb': RGBColor(0xe5, 0x4b, 0x4b), 'label': 'PPAP'},
    'Top-level initiative': {'rgb': RGBColor(0x3b, 0x5b, 0xdb), 'label': 'Init'},
}
DEFAULT_TYPE = {'rgb': RGBColor(0x82, 0x82, 0x82), 'label': '?'}
COMPLETED_COLOR = RGBColor(0x27, 0xae, 0x60)

# Layout constants
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)
LEFT_MARGIN = Inches(0.3)
TOP_MARGIN = Inches(0.6)
HEADER_HEIGHT = Inches(0.3)
ROW_HEIGHT = Inches(0.3)
LABEL_WIDTH = Inches(3.2)
CHART_LEFT = LEFT_MARGIN + LABEL_WIDTH + Inches(0.1)
CHART_WIDTH = SLIDE_WIDTH - CHART_LEFT - Inches(0.3)
TASKS_PER_SLIDE = 18


# ========================================
# Date Parsing
# ========================================
def parse_date(date_str):
    if not date_str or not date_str.strip():
        return None
    s = date_str.strip()
    # Remove time portions: "12:00 午前", "12:00 AM", etc.
    s = re.sub(r'\s+\d{1,2}:\d{2}(:\d{2})?\s*(午前|午後|AM|PM)?', '', s, flags=re.IGNORECASE).strip()
    
    # Format: DD/M/YY  (Jira: day first, 2-digit year last, e.g. "20/4/26" = Apr 20 2026)
    m = re.match(r'^(\d{1,4})/(\d{1,2})/(\d{1,4})$', s)
    if m:
        a, b, c = int(m.group(1)), int(m.group(2)), int(m.group(3))
        if a > 31:
            # YYYY/M/D
            year, month, day = a, b, c
        elif c <= 99:
            # DD/M/YY — Jira format
            day, month, year = a, b, c + 2000
        else:
            # D/M/YYYY
            day, month, year = a, b, c
        try:
            return datetime(year, month, day)
        except ValueError:
            return None
    
    # Format: YYYY-MM-DD
    m = re.match(r'^(\d{4})-(\d{1,2})-(\d{1,2})$', s)
    if m:
        try:
            return datetime(int(m.group(1)), int(m.group(2)), int(m.group(3)))
        except ValueError:
            return None
    return None



def parse_iso_date(s):
    """Parse ISO date string from JSON (YYYY-MM-DD or ISO 8601)."""
    if not s:
        return None
    try:
        # Handle "YYYY/MM/DD" and "YYYY-MM-DD"
        s = s.strip().replace('/', '-')
        # Remove time portion if present
        if 'T' in s:
            s = s.split('T')[0]
        parts = s.split('-')
        if len(parts) == 3:
            return datetime(int(parts[0]), int(parts[1]), int(parts[2]))
    except (ValueError, IndexError):
        pass
    return None


# ========================================
# Timeline Helpers
# ========================================
def days_in_month(dt):
    if dt.month == 12:
        return (datetime(dt.year + 1, 1, 1) - datetime(dt.year, dt.month, 1)).days
    return (datetime(dt.year, dt.month + 1, 1) - datetime(dt.year, dt.month, 1)).days


def get_months_between(start, end):
    months = []
    current = datetime(start.year, start.month, 1)
    while current < end:
        months.append(current)
        if current.month == 12:
            current = datetime(current.year + 1, 1, 1)
        else:
            current = datetime(current.year, current.month + 1, 1)
    return months


def get_mondays_between(start, end):
    """Return a list of Mondays covering [start, end).
    First Monday is on-or-before start; last Monday's week still overlaps end.
    """
    aligned = start - timedelta(days=start.weekday())
    mondays = []
    current = aligned
    while current < end:
        mondays.append(current)
        current += timedelta(days=7)
    if not mondays:
        mondays.append(aligned)
    return mondays


# ========================================
# Helpers
# ========================================
def add_text_box(slide, text, x, y, w, h,
                 font_size=7, color=TEXT_PRIMARY, bold=False,
                 align=PP_ALIGN.LEFT, font_name='Meiryo UI'):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    tf._txBody.bodyPr.set('anchor', 'ctr')
    return txBox


# ========================================
# PPTX Generation from JSON data
# ========================================
def generate_pptx_from_data(tasks_data, output_filename, title="Gantt Chart", timeline_start=None, timeline_end=None):
    """Generate PPTX with Smartsheet-style weekly columns, no type grouping,
    and dynamic row height that scales with the timeline range.
    """

    # ---- Parse task dates ----
    tasks = []
    for t in tasks_data:
        start = parse_iso_date(t.get('startDate', ''))
        end = parse_iso_date(t.get('endDate', ''))
        if start or end:
            tasks.append({
                'key': t.get('key', ''),
                'type': t.get('type', 'Task'),
                'summary': t.get('summary', ''),
                'status': t.get('status', 'To Do'),
                'assignee': t.get('assignee', ''),
                'start_date': start,
                'end_date': end,
                'parent': t.get('parent', ''),
            })

    if not tasks:
        return None, "No tasks with dates found"

    # ---- Flat sort by start date (no grouping) ----
    tasks.sort(key=lambda t: t['start_date'] or t['end_date'] or datetime(9999, 1, 1))

    # ---- Timeline bounds, Monday-aligned ----
    custom_start = parse_iso_date(timeline_start) if timeline_start else None
    custom_end = parse_iso_date(timeline_end) if timeline_end else None

    all_dates = []
    for t in tasks:
        if t['start_date']: all_dates.append(t['start_date'])
        if t['end_date']: all_dates.append(t['end_date'])

    earliest = min(all_dates)
    latest = max(all_dates)

    raw_start = custom_start if custom_start else earliest - timedelta(days=14)
    raw_end = (custom_end if custom_end else latest) + timedelta(days=14)

    mondays = get_mondays_between(raw_start, raw_end)
    if not mondays:
        return None, "Invalid timeline range"

    bounds_start = mondays[0]
    bounds_end = mondays[-1] + timedelta(days=7)
    total_days = (bounds_end - bounds_start).days
    total_weeks = len(mondays)
    today = datetime.now()

    # ---- Two-row header heights (Smartsheet style) ----
    MONTH_HEADER_H = Inches(0.28)
    WEEK_HEADER_H = Inches(0.22)
    HEADER_TOTAL_H = MONTH_HEADER_H + WEEK_HEADER_H
    LEGEND_H = Inches(0.4)

    # ---- Dynamic row height — scales with timeline zoom ----
    available_body_h = SLIDE_HEIGHT - TOP_MARGIN - HEADER_TOTAL_H - LEGEND_H - Inches(0.15)
    MIN_ROW = Inches(0.22)
    MAX_ROW = Inches(0.40)
    n_tasks = len(tasks)

    ideal_rh = available_body_h // max(n_tasks, 1)
    if ideal_rh >= MIN_ROW:
        row_height = min(ideal_rh, MAX_ROW)
        tasks_per_slide = n_tasks
    else:
        row_height = MIN_ROW
        tasks_per_slide = max(int(available_body_h // row_height), 1)

    bar_h = int(row_height * 0.6)

    # ---- Paginate (flat list) ----
    pages = [tasks[i:i + tasks_per_slide] for i in range(0, n_tasks, tasks_per_slide)]

    # ---- Chart dimensions (equal-width weekly columns) ----
    chart_width_emu = CHART_WIDTH
    week_width = chart_width_emu / total_weeks
    day_width = week_width / 7

    # ---- Create presentation ----
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    month_names_jp = ['1月','2月','3月','4月','5月','6月','7月','8月','9月','10月','11月','12月']

    for page_idx, page_tasks in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR

        # Title
        page_label = f" ({page_idx+1}/{len(pages)})" if len(pages) > 1 else ""
        add_text_box(slide, f"{title}{page_label}",
                     LEFT_MARGIN, Inches(0.12), Inches(10), Inches(0.4),
                     font_size=16, color=TITLE_COLOR, bold=True)

        add_text_box(slide,
                     f"{bounds_start.strftime('%Y/%m/%d')} - {bounds_end.strftime('%Y/%m/%d')}  ({total_weeks} weeks)",
                     LEFT_MARGIN, Inches(0.42), Inches(8), Inches(0.18),
                     font_size=7, color=TEXT_MUTED)

        # ---- Top header row: month spans ----
        cur_ym = None
        span_start = 0
        for i in range(len(mondays) + 1):
            ym = (mondays[i].year, mondays[i].month) if i < len(mondays) else None
            if cur_ym is None and ym is not None:
                cur_ym = ym
                span_start = i
            elif ym != cur_ym:
                x = CHART_LEFT + int(week_width * span_start)
                w = int(week_width * i) - int(week_width * span_start)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, x, TOP_MARGIN, w, MONTH_HEADER_H
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = BG_SECONDARY
                shape.line.color.rgb = GRID_COLOR
                shape.line.width = Pt(0.5)
                tf = shape.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.text = f"{cur_ym[0]}/{month_names_jp[cur_ym[1] - 1]}"
                p.font.size = Pt(9)
                p.font.color.rgb = TITLE_COLOR
                p.font.bold = True
                p.font.name = 'Meiryo UI'
                p.alignment = PP_ALIGN.CENTER
                tf._txBody.bodyPr.set('anchor', 'ctr')
                cur_ym = ym
                span_start = i

        # ---- Bottom header row: Monday date per week ----
        week_hdr_y = TOP_MARGIN + MONTH_HEADER_H
        for i, mon in enumerate(mondays):
            x = CHART_LEFT + int(week_width * i)
            w = int(week_width * (i + 1)) - int(week_width * i)
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, week_hdr_y, w, WEEK_HEADER_H
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_COLOR
            cell.line.color.rgb = GRID_COLOR
            cell.line.width = Pt(0.25)
            tf = cell.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"{mon.month}/{mon.day}"
            first_of_month = mon.day <= 7
            p.font.size = Pt(7)
            p.font.color.rgb = TEXT_PRIMARY if first_of_month else TEXT_SECONDARY
            p.font.bold = first_of_month
            p.font.name = 'Meiryo UI'
            p.alignment = PP_ALIGN.CENTER
            tf._txBody.bodyPr.set('anchor', 'ctr')

        # ---- Body geometry ----
        body_y = TOP_MARGIN + HEADER_TOTAL_H
        body_h = int(row_height * len(page_tasks))

        # Weekly grid lines
        for i in range(total_weeks + 1):
            x = CHART_LEFT + int(week_width * i)
            gl = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, body_y, Pt(0.5), body_h
            )
            gl.fill.solid()
            gl.fill.fore_color.rgb = GRID_COLOR
            gl.line.fill.background()

        # Today line
        today_days = (today - bounds_start).days
        if 0 <= today_days <= total_days:
            today_x = int(CHART_LEFT + day_width * today_days)
            tl = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, today_x, TOP_MARGIN - Inches(0.05),
                Pt(2), HEADER_TOTAL_H + body_h + Inches(0.1)
            )
            tl.fill.solid()
            tl.fill.fore_color.rgb = TODAY_COLOR
            tl.line.fill.background()
            add_text_box(slide, 'Today',
                         today_x - Inches(0.2), TOP_MARGIN - Inches(0.22),
                         Inches(0.4), Inches(0.16),
                         font_size=6, color=TODAY_COLOR, align=PP_ALIGN.CENTER, bold=True)

        # ---- Task rows (flat, no grouping) ----
        for row_idx, item in enumerate(page_tasks):
            y = body_y + int(row_idx * row_height)

            if row_idx % 2 == 0:
                rb = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, LEFT_MARGIN, y,
                    SLIDE_WIDTH - LEFT_MARGIN * 2, row_height
                )
                rb.fill.solid()
                rb.fill.fore_color.rgb = BG_ROW_ALT
                rb.line.fill.background()

            tc = TYPE_COLORS.get(item['type'], DEFAULT_TYPE)

            add_text_box(slide, item['key'],
                         LEFT_MARGIN + Inches(0.05), y,
                         Inches(0.7), row_height,
                         font_size=6, color=TEXT_MUTED)

            badge_w = Inches(0.4)
            badge_h_actual = min(bar_h, Inches(0.2))
            badge_y = y + (row_height - badge_h_actual) // 2
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                LEFT_MARGIN + Inches(0.75), badge_y, badge_w, badge_h_actual
            )
            badge.fill.solid()
            r, g, b = tc['rgb'][0], tc['rgb'][1], tc['rgb'][2]
            badge.fill.fore_color.rgb = RGBColor(
                min(255, int(r * 0.15 + 255 * 0.85)),
                min(255, int(g * 0.15 + 255 * 0.85)),
                min(255, int(b * 0.15 + 255 * 0.85))
            )
            badge.line.fill.background()
            tf = badge.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = tc['label']
            p.font.size = Pt(5)
            p.font.color.rgb = tc['rgb']
            p.font.bold = True
            p.font.name = 'Meiryo UI'
            p.alignment = PP_ALIGN.CENTER
            tf._txBody.bodyPr.set('anchor', 'ctr')

            add_text_box(slide, item['summary'],
                         LEFT_MARGIN + Inches(1.2), y,
                         LABEL_WIDTH - Inches(1.2), row_height,
                         font_size=7, color=TEXT_PRIMARY)

            if item['start_date'] or item['end_date']:
                start = item['start_date'] or item['end_date']
                end = item['end_date'] or item['start_date']

                start_days = (start - bounds_start).days
                duration_days = max((end - start).days, 1)

                bar_x = int(CHART_LEFT + day_width * start_days)
                bar_w = max(int(day_width * duration_days), Inches(0.06))
                bar_y = y + (row_height - bar_h) // 2

                is_completed = item['status'] in ('完了', 'Done')
                bar_color = COMPLETED_COLOR if is_completed else tc['rgb']

                if item['type'] == 'Milestone':
                    dia = min(bar_h, Inches(0.2))
                    diamond = slide.shapes.add_shape(
                        MSO_SHAPE.DIAMOND,
                        bar_x - dia // 2, y + (row_height - dia) // 2,
                        dia, dia
                    )
                    diamond.fill.solid()
                    diamond.fill.fore_color.rgb = bar_color
                    diamond.line.fill.background()
                else:
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        bar_x, bar_y, bar_w, bar_h
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = bar_color
                    bar.line.fill.background()

                    if bar_w > Inches(0.8):
                        tf = bar.text_frame
                        tf.word_wrap = False
                        p = tf.paragraphs[0]
                        p.text = item['summary']
                        p.font.size = Pt(6 if row_height < Inches(0.3) else 7)
                        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                        p.font.name = 'Meiryo UI'
                        p.alignment = PP_ALIGN.LEFT
                        tf._txBody.bodyPr.set('anchor', 'ctr')

                    if start_days >= 0:
                        add_text_box(slide, start.strftime('%m/%d'),
                                     bar_x - Inches(0.35), bar_y,
                                     Inches(0.33), bar_h,
                                     font_size=4, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

                    if duration_days > 1:
                        add_text_box(slide, end.strftime('%m/%d'),
                                     bar_x + bar_w + Inches(0.02), bar_y,
                                     Inches(0.33), bar_h,
                                     font_size=4, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

        # ---- Legend ----
        legend_y = SLIDE_HEIGHT - Inches(0.35)
        legend_x = LEFT_MARGIN

        for label, color in [
            ('Milestone', TYPE_COLORS['Milestone']['rgb']),
            ('Goal', TYPE_COLORS['Goal']['rgb']),
            ('Review', TYPE_COLORS['Review']['rgb']),
            ('Event', TYPE_COLORS['Event']['rgb']),
            ('Task', TYPE_COLORS['Task']['rgb']),
            ('Validation', TYPE_COLORS['Validation']['rgb']),
            ('Completed', COMPLETED_COLOR),
        ]:
            sw = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                legend_x, legend_y + Inches(0.03), Inches(0.2), Inches(0.1)
            )
            sw.fill.solid()
            sw.fill.fore_color.rgb = color
            sw.line.fill.background()

            add_text_box(slide, label,
                         legend_x + Inches(0.22), legend_y,
                         Inches(0.6), Inches(0.15),
                         font_size=6, color=TEXT_SECONDARY)
            legend_x += Inches(0.9)

        ts = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_x, legend_y, Pt(2), Inches(0.15))
        ts.fill.solid()
        ts.fill.fore_color.rgb = TODAY_COLOR
        ts.line.fill.background()
        add_text_box(slide, 'Today', legend_x + Inches(0.05), legend_y, Inches(0.5), Inches(0.15),
                     font_size=6, color=TODAY_COLOR)

    # Save
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / output_filename
    prs.save(str(output_path))
    return str(output_path.resolve()), None


# ========================================
# HTTP Handler
# ========================================
class GanttHandler(http.server.SimpleHTTPRequestHandler):
    
    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=str(SCRIPT_DIR), **kwargs)
    
    def do_POST(self):
        if self.path == '/api/generate-pptx':
            content_length = int(self.headers.get('Content-Length', 0))
            body = self.rfile.read(content_length)
            
            try:
                data = json.loads(body.decode('utf-8'))
                tasks = data.get('tasks', [])
                title = data.get('title', 'Gantt Chart')
                filename = data.get('filename', f'gantt_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx')
                
                # Sanitize filename
                filename = re.sub(r'[^\w\-_\.]', '_', filename)
                if not filename.endswith('.pptx'):
                    filename += '.pptx'
                
                # Extract custom timeline bounds if provided
                tl_start = data.get('timelineStart', None)
                tl_end = data.get('timelineEnd', None)
                
                output_path, error = generate_pptx_from_data(
                    tasks, filename, title,
                    timeline_start=tl_start,
                    timeline_end=tl_end
                )
                
                if error:
                    self.send_response(400)
                    self.send_header('Content-Type', 'application/json')
                    self.end_headers()
                    self.wfile.write(json.dumps({'error': error}).encode('utf-8'))
                else:
                    self.send_response(200)
                    self.send_header('Content-Type', 'application/json')
                    self.end_headers()
                    result = {
                        'success': True,
                        'path': output_path,
                        'filename': filename,
                    }
                    self.wfile.write(json.dumps(result).encode('utf-8'))
                    print(f"  [OK] Generated: {output_path}")
                    
            except Exception as e:
                self.send_response(500)
                self.send_header('Content-Type', 'application/json')
                self.end_headers()
                self.wfile.write(json.dumps({'error': str(e)}).encode('utf-8'))
                print(f"  [ERROR] {e}")
        else:
            self.send_response(404)
            self.end_headers()
    
    def log_message(self, format, *args):
        # Quieter logging
        if '/api/' in str(args[0]) if args else False:
            print(f"  [API] {args[0]}")
        else:
            super().log_message(format, *args)


# ========================================
# Main
# ========================================
if __name__ == '__main__':
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    print("=" * 50)
    print("  Gantt Chart Generator Server")
    print("=" * 50)
    print(f"  URL:    http://localhost:{PORT}")
    print(f"  Output: {OUTPUT_DIR.resolve()}")
    print(f"  Press Ctrl+C to stop")
    print("=" * 50)
    
    server = http.server.HTTPServer(('', PORT), GanttHandler)
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\n  Server stopped.")
        server.server_close()
